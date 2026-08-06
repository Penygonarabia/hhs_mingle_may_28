from odoo import http
from odoo.http import request
from .sales_mail_main import PbiSalesMailController
from .sales_kpi_main import PbiSalesKpiController, CUSTOMER_NAME_EXPR, _FRANCHISE_SCOPE_SQL, _FRANCHISE_SCOPE_PARAMS, AMOUNT_EXPR

SALESMAN_EXPR = """
    (SELECT rp.name FROM res_users ru
     JOIN res_partner rp ON rp.id = ru.partner_id
     WHERE ru.is_salesman = true AND ru.user_code = bi_salesmancode
     LIMIT 1)
"""


class PbiSalesKpiQuarterlyController(PbiSalesKpiController):
    """Backs "Sales Dashboard (Q)" under "Quarterly Sales Dashboards".
    Subclasses the Sales KPI controller to target transaction tables via
    v_pbi_sales_analysis_q and customizes levels to use Sales Type Group/Sub-group
    and include Salesman level.
    """

    LEVELS = ["salesTypeGroup", "salesTypeSubGroup", "region", "salesman", "customer", "productGroup", "productSubGroup"]

    LEVEL_COLUMNS = {
        "salesTypeGroup": ("bi_csttypecode", "bi_csttypedesc"),
        "salesTypeSubGroup": ("bi_cstsubtypecode", "bi_cstsubtypedesc"),
        "region": ("bi_cstregioncode", "bi_cstregiondesc"),
        "salesman": ("bi_salesmancode", f"({SALESMAN_EXPR})"),
        "customer": ("bi_cstno", f"({CUSTOMER_NAME_EXPR})"),
        "productGroup": ("bi_pgroupcode", "bi_pgroupname"),
        "productSubGroup": ("bi_psgroupcode", "bi_psgroupname"),
    }

    LEVEL_LABELS = {
        "salesTypeGroup": "Sales Type Group",
        "salesTypeSubGroup": "Sales Type Sub-Group",
        "region": "Region",
        "salesman": "Salesman",
        "customer": "Customer",
        "productGroup": "Product Group",
        "productSubGroup": "Product Sub-Group",
    }

    def _query(self, sql, params=()):
        # Rewrite query to point to the transaction-based view instead of bidata
        sql = sql.replace("FROM bidata", "FROM v_pbi_sales_analysis_q")
        sql = sql.replace("bidata.bi_", "v_pbi_sales_analysis_q.bi_")
        sql = sql.replace(" bidata ", " v_pbi_sales_analysis_q ")
        return super()._query(sql, params)

    def _dim_clauses(self, franchise, customer_group, sales_type_group):
        clauses, params = [], []
        if franchise and franchise != "all":
            clauses.append("bi_franchisename = %s")
            params.append(franchise)
        if sales_type_group and sales_type_group != "all":
            # For Q dashboard, sales_type_group is the actual salgrp_ref code (e.g. RT, WS, CP)
            clauses.append("bi_csttypecode = %s")
            params.append(sales_type_group)
        if customer_group and customer_group != "all":
            clauses.append("bi_cstno IN (SELECT customer_no FROM v_customergroups WHERE customer_groupcode = %s)")
            params.append(customer_group)
        return clauses, params

    def _breakdown(self, level, year, franchise, customer_group, sales_type_group, drill_path, month=None, month_lte=None):
        # If we are on salesTypeGroup, we don't want the parent's fixed-mapping code because it expects bidata customer types.
        # We can just run the query and return the list of rows normally, or fetch all groups from salestypes_group to zero-fill.
        if level == "salesTypeGroup" and not drill_path and (not sales_type_group or sales_type_group == "all"):
            code_col, name_col = self.LEVEL_COLUMNS[level]
            clauses, params = self._period_clauses(year, month, month_lte)
            dim_clauses, dim_params = self._dim_clauses(franchise, customer_group, sales_type_group)
            drill_clauses, drill_params = self._drill_clauses(drill_path)
            
            # Fetch all groups from salestypes_group so we can zero-fill
            groups = self._query("SELECT salgrp_ref AS code, salgrp_name AS label FROM salestypes_group")
            
            # Run normal breakdown query
            clauses = clauses + [_FRANCHISE_SCOPE_SQL] + dim_clauses + drill_clauses
            params = params + _FRANCHISE_SCOPE_PARAMS + dim_params + drill_params
            where = " AND ".join(clauses)
            rows = self._query(f"""
                SELECT {code_col} AS code, {name_col} AS label, sum({AMOUNT_EXPR}) AS sales, sum(bi_budgetamount) AS budget,
                       sum(bi_qty) AS qty, sum(bi_budgetqty) AS budget_qty
                FROM v_pbi_sales_analysis_q WHERE {where}
                GROUP BY 1, 2
                ORDER BY 3 DESC
            """, params)
            
            by_code = {r["code"]: r for r in rows}
            out = []
            for g in groups:
                code, label = g["code"], g["label"]
                r = by_code.get(code)
                out.append({
                    "code": code, "label": label,
                    "sales": float(r["sales"] or 0) if r else 0.0,
                    "budget": float(r["budget"] or 0) if r else 0.0,
                    "qty": int(r["qty"] or 0) if r else 0,
                    "budgetQty": int(r["budget_qty"] or 0) if r else 0,
                })
            out.sort(key=lambda o: o["sales"], reverse=True)
            return out

        return super()._breakdown(level, year, franchise, customer_group, sales_type_group, drill_path, month, month_lte)

    def _has_access_q(self):
        return self._has_access_for("pbi_dashboards.menu_pbi_sales_kpi_analysis_q")

    @http.route('/pbi_dashboards/sales_kpi_q/data', type='json', auth='user')
    def sales_kpi_data_q(self, period=None, franchise='all', customerGroup='all', salesTypeGroup='all', drillPath=None, levelFilterCode=None, viewLevel=None):
        return self._sales_kpi_data(self._has_access_q(), period, franchise, customerGroup, salesTypeGroup, drillPath, levelFilterCode, viewLevel)


class PbiSalesQuarterlyController(PbiSalesMailController):
    """Backs "Sales Analysis (Q)" dashboard under "Quarterly Sales Dashboards".
    Subclasses the Sales Mail controller to reuse all formatting, narrative generation,
    and PPTX export/notes-import features, but overrides _query to rewrite all queries
    to target the transaction tables view v_pbi_sales_analysis_q instead of bidata.
    """

    def _query(self, sql, params=()):
        # Rewrite query to point to the transaction-based view instead of bidata
        sql = sql.replace("FROM bidata", "FROM v_pbi_sales_analysis_q")
        sql = sql.replace("bidata.bi_", "v_pbi_sales_analysis_q.bi_")
        sql = sql.replace(" bidata ", " v_pbi_sales_analysis_q ")
        return super()._query(sql, params)

    def _build_pptx(self, year, month, franchise):
        # Build base PPTX presentation
        data = super()._build_pptx(year, month, franchise)
        
        # Intercept and replace "Source: bidata" label with transaction source description
        from pptx import Presentation
        from io import BytesIO
        
        prs = Presentation(BytesIO(data))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if "Source: bidata" in paragraph.text:
                            paragraph.text = paragraph.text.replace(
                                "Source: bidata", 
                                "Source: transaction_header / transaction_details"
                            )
        
        stream = BytesIO()
        prs.save(stream)
        return stream.getvalue()

    def _has_access_q(self):
        return self._has_access_for("pbi_dashboards.menu_pbi_sales_analysis_q")

    @http.route('/pbi_dashboards/sales_analysis_q/data', type='json', auth='user')
    def sales_analysis_q_data(self, period=None, franchise='Midea'):
        return self._sales_mail_data(self._has_access_q(), period, franchise)

    @http.route('/pbi_dashboards/sales_analysis_q/export.pptx', type='http', auth='user')
    def export_pptx_q(self, period=None, franchise='Midea', **kwargs):
        return self._export_pptx(self._has_access_q(), period, franchise)

    @http.route('/pbi_dashboards/sales_analysis_q/import_notes', type='http', auth='user', methods=['POST'])
    def import_notes_q(self, period=None, franchise='Midea', **kwargs):
        return self._import_notes(self._has_access_q(), period, franchise)
