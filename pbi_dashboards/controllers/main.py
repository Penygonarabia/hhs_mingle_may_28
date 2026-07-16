import calendar
import json
import re
from io import BytesIO

from odoo import http, fields
from odoo.http import request

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

MONTH_NAMES = ["January", "February", "March", "April", "May", "June",
               "July", "August", "September", "October", "November", "December"]

# The "Year" filter also accepts a single "YYYY-MM" month value (e.g. "2025-03").
MONTH_VALUE_RE = re.compile(r"^(\d{4})-(\d{2})$")

PPTX_COLOR_2024 = RGBColor(0x2a, 0x78, 0xd6)
PPTX_COLOR_2025 = RGBColor(0x1b, 0xaf, 0x7a)
PPTX_COLOR_BUDGET = RGBColor(0xed, 0xa1, 0x00)

# Maps a chart/table slide's title text (substring match, case-insensitive)
# back to the on-screen narrative section it corresponds to — lets the notes
# importer match slides by content rather than fixed position, so it still
# works if slides get reordered in PowerPoint. Order mirrors _build_pptx.
PPTX_NOTE_SECTION_MATCHERS = (
    ("monthly sales", "trend"),
    ("sales by region", "region"),
    ("sales by franchise", "franchise"),
    ("sales by customer type", "customerType"),
    ("top product groups", "product"),
    ("top 10 salesmen", "salesman"),
)


def _pptx_slide_title(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape.text_frame.text.strip()
    return ""


def _pptx_note_section_key(title):
    lowered = title.lower()
    for needle, key in PPTX_NOTE_SECTION_MATCHERS:
        if needle in lowered:
            return key
    return None


# Mirrors the JS-side NOTE_HEADING_RE / highlightFigures in sales_dashboard.js
# so a slide's speaker notes are bold-formatted the same way the on-screen
# narrative is: the YTD line as a heading, figures (M/K amounts, comma-grouped
# numbers, percentages) in bold within ordinary sentences.
PPTX_NOTES_HEADING_RE = re.compile(r'^Year-to-date .+ performance$', re.IGNORECASE)
PPTX_NOTES_FIGURE_RE = re.compile(r'[+-]?\d{1,3}(?:,\d{3})+(?:\.\d+)?\b|[+-]?\d+(?:\.\d+)?[MK]\b|[+-]?\d+(?:\.\d+)?%')
PPTX_NOTES_BOLD_RUN_RE = re.compile(r'\*\*(.+?)\*\*')


def _pptx_marked_text(plain_text):
    """Turns a freshly-generated plain-text narrative into the canonical
    **bold**/##heading## marked form — the same bold spans _pptx_write_notes
    turns into real runs, and what an untouched re-import's actual run
    formatting (see _pptx_read_notes_marked) is compared against to tell a
    genuine hand edit apart from notes the user left alone."""
    lines = []
    for line in plain_text.split('\n'):
        if PPTX_NOTES_HEADING_RE.match(line):
            lines.append(f"##{line}##")
            continue
        pos = 0
        out = []
        for m in PPTX_NOTES_FIGURE_RE.finditer(line):
            if m.start() > pos:
                out.append(line[pos:m.start()])
            out.append(f"**{m.group(0)}**")
            pos = m.end()
        out.append(line[pos:])
        lines.append(''.join(out))
    return '\n'.join(lines)


def _pptx_fill_marked_text(text_frame, plain_text, heading_size=Pt(14), body_size=None, already_marked=False):
    """Shared renderer for a **bold**/##heading## marked narrative into any
    pptx text frame — PowerPoint's hidden Notes pane (_pptx_write_notes) or
    a visible on-slide textbox (the sales-mail dashboards' notes sidebar,
    mirroring the web page's `.narrative` column) both go through this.
    `already_marked=True` skips the auto-figure-marking pass for text that
    is already in canonical **bold**/##heading## form — e.g. a hand-edited
    note round-tripped back from pbi.dashboard.note via
    _pptx_read_notes_marked, which already carries its own (possibly
    non-figure, user-added) bold spans. Re-running _pptx_marked_text on
    that would scan for numeric figures inside text that already has
    literal '**' in it, doubling up the markers and corrupting every bold
    span from the first stray match onward."""
    text_frame.clear()
    marked = plain_text if already_marked else _pptx_marked_text(plain_text)
    for i, line in enumerate(marked.split('\n')):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.space_after = Pt(10)
        heading_m = re.fullmatch(r'##(.*)##', line)
        if heading_m:
            run = p.add_run()
            run.text = heading_m.group(1)
            run.font.bold = True
            run.font.size = heading_size
            continue
        pos = 0
        for m in PPTX_NOTES_BOLD_RUN_RE.finditer(line):
            if m.start() > pos:
                r = p.add_run()
                r.text = line[pos:m.start()]
                if body_size:
                    r.font.size = body_size
            run = p.add_run()
            run.text = m.group(1)
            run.font.bold = True
            if body_size:
                run.font.size = body_size
            pos = m.end()
        if pos < len(line):
            r = p.add_run()
            r.text = line[pos:]
            if body_size:
                r.font.size = body_size


def _pptx_write_notes(slide, plain_text):
    _pptx_fill_marked_text(slide.notes_slide.notes_text_frame, plain_text)


def _pptx_read_notes_marked(text_frame):
    """The import-time counterpart to _pptx_marked_text: reconstructs the
    **bold**/##heading## marked form from a notes text frame's ACTUAL run
    formatting, so any bold the user added or removed by hand in PowerPoint
    (not just the figures/heading auto-bolded at export time) round-trips
    faithfully into the stored override and the on-screen rendering."""
    lines = []
    for p in text_frame.paragraphs:
        line = ''.join(f"**{run.text}**" if run.font.bold else run.text for run in p.runs)
        plain_line = line.replace('**', '')
        if PPTX_NOTES_HEADING_RE.match(plain_line):
            line = f"##{plain_line}##"
        lines.append(line)
    return '\n'.join(lines)


def _pfmt_m(n):
    return f"{(n or 0) / 1e6:.1f}M"


def _pfmt(n):
    return f"{int(round(n or 0)):,}"


def _pshare(part, total):
    return f"{round((part or 0) / total * 100)}%" if total else "–"


def _pnarrative_kpis_single(kpis):
    achv = (kpis['sales'] / kpis['budget'] * 100) if kpis.get('budget') else None
    partial = bool(kpis.get('isPartialMonth'))
    bits = [f"Total sales {_pfmt_m(kpis['sales'])}",
            f"{'target' if partial else 'total budget'} {_pfmt_m(kpis['budget'])}"]
    if kpis.get('prevYearSales') is not None:
        bits.append(f"last year {_pfmt_m(kpis['prevYearSales'])}")
    if achv is not None:
        bits.append(f"achievement {achv:.1f}% vs {'full-month budget' if partial else 'budget'}")
    bits.append(f"{_pfmt(kpis['qty'])} units sold")
    bits.append(f"{_pfmt(kpis['customers'])} active customers")
    text = ", ".join(bits) + "."
    if partial and kpis.get('asOfDate'):
        as_of = fields.Date.from_string(kpis['asOfDate'][:10])
        text += (f" Data is only synced through {as_of.strftime('%B %d, %Y')} for this month"
                 " — sales figures are month-to-date, while the budget shown is the full-month target.")
    return text


def _pnarrative_kpis_compare(k24, k25):
    delta = lambda a, b: ((b - a) / a * 100) if a else None
    sd, bd, cd = delta(k24['sales'], k25['sales']), delta(k24['budget'], k25['budget']), delta(k24['customers'], k25['customers'])
    achv = (k25['sales'] / k25['budget'] * 100) if k25.get('budget') else None
    bits = [f"2025 sales {_pfmt_m(k25['sales'])}" + (f" ({sd:+.1f}% vs 2024)" if sd is not None else "")]
    bits.append(f"2025 budget {_pfmt_m(k25['budget'])}" + (f" ({bd:+.1f}% vs 2024)" if bd is not None else ""))
    if achv is not None:
        bits.append(f"achievement {achv:.1f}% vs 2025 budget")
    bits.append(f"{_pfmt(k25['customers'])} active customers" + (f" ({cd:+.1f}% vs 2024)" if cd is not None else ""))
    return ", ".join(bits) + "."


def _pfranchise_label(franchise):
    return 'the selected franchises' if franchise == 'all' else franchise


# The following four mirror monthValueSentence / monthQtySentence /
# ytdValueSentence / ytdQtySentence in sales_dashboard.js sentence-for-sentence
# (plain text, no <b>/citation markup) — kept as separate functions, same as
# the JS, rather than folded into _pnarrative_trend.
def _pmonth_value_sentence(month_label, f_label, cur, prev):
    has_budget = (cur.get('budget') or 0) > 0
    has_prev = bool(prev)
    s = f"For {month_label}, {f_label} achieved sales value of {_pfmt_m(cur.get('sales'))}"
    if has_budget and has_prev:
        verb = 'outperforming' if cur['sales'] >= cur['budget'] else 'underperforming'
        cmp_word = 'still above' if cur['sales'] >= prev['sales'] else 'below'
        s += (f" versus a target of {_pfmt_m(cur['budget'])}, {verb} the value target, "
              f"though this was {cmp_word} last year's {_pfmt_m(prev['sales'])} for the same period")
    elif has_budget:
        verb = 'outperforming' if cur['sales'] >= cur['budget'] else 'underperforming'
        s += f" versus a target of {_pfmt_m(cur['budget'])}, {verb} the value target"
    elif has_prev:
        cmp_word = 'above' if cur['sales'] >= prev['sales'] else 'below'
        s += f", {cmp_word} last year's {_pfmt_m(prev['sales'])} for the same period"
    return s


def _pmonth_qty_sentence(month_label, cur, prev):
    has_budget = (cur.get('budgetQty') or 0) > 0
    has_prev = bool(prev)
    s = f"In volume terms, {month_label} sales reached {_pfmt(cur.get('qty'))} units"
    if has_budget and has_prev:
        growth = 'growth' if cur['qty'] >= prev['qty'] else 'a decline'
        target_phrase = ('and ahead of the volume target' if cur['qty'] >= cur['budgetQty']
                          else 'but slightly below the volume target')
        s += (f" compared with a target of {_pfmt(cur['budgetQty'])} units and {_pfmt(prev['qty'])} units last year, "
              f"indicating {growth} over last year {target_phrase}")
    elif has_budget:
        cmp_word = 'ahead of' if cur['qty'] >= cur['budgetQty'] else 'below'
        s += f" compared with a target of {_pfmt(cur['budgetQty'])} units, {cmp_word} the volume target"
    elif has_prev:
        growth = 'growth' if cur['qty'] >= prev['qty'] else 'a decline'
        s += f" compared with {_pfmt(prev['qty'])} units last year, indicating {growth} over last year"
    return s


def _pytd_value_sentence(year_label, f_label, cur, prev):
    has_budget = (cur.get('budget') or 0) > 0
    has_prev = bool(prev)
    s = f"For {year_label} year-to-date, {f_label}'s sales value is {_pfmt_m(cur.get('sales'))}"
    if has_budget and has_prev:
        gap_pct = abs(cur['sales'] - cur['budget']) / cur['budget'] * 100
        closeness = 'narrowly ' if gap_pct <= 5 else ''
        target_cmp = f"{closeness}above" if cur['sales'] >= cur['budget'] else f"{closeness}below"
        ahead = cur['sales'] >= prev['sales']
        s += (f", {target_cmp} the target of {_pfmt_m(cur['budget'])}, but clearly "
              f"{'ahead of' if ahead else 'behind'} last year's {_pfmt_m(prev['sales'])}, "
              f"reflecting {'a solid uplift' if ahead else 'a shortfall'} versus the prior year")
    elif has_budget:
        s += f", {'above' if cur['sales'] >= cur['budget'] else 'below'} the target of {_pfmt_m(cur['budget'])}"
    elif has_prev:
        s += f", {'ahead of' if cur['sales'] >= prev['sales'] else 'behind'} last year's {_pfmt_m(prev['sales'])}"
    return s


def _pytd_qty_sentence(year_label, cur, prev):
    has_budget = (cur.get('budgetQty') or 0) > 0
    has_prev = bool(prev)
    s = f"In terms of quantity, {year_label} YTD sales total {_pfmt(cur.get('qty'))} units"
    if has_budget and has_prev:
        target_cmp = 'above the target of' if cur['qty'] >= cur['budgetQty'] else 'below the target of'
        prev_cmp = "above last year's" if cur['qty'] >= prev['qty'] else "below last year's"
        plan_word = 'ahead of plan' if cur['qty'] >= cur['budgetQty'] else 'behind plan'
        yoy_word = 'improving year on year' if cur['qty'] >= prev['qty'] else 'declining year on year'
        s += (f", which is {target_cmp} {_pfmt(cur['budgetQty'])} units and {prev_cmp} {_pfmt(prev['qty'])} units, "
              f"showing that overall unit performance is {plan_word} and {yoy_word}")
    elif has_budget:
        s += f", {'above' if cur['qty'] >= cur['budgetQty'] else 'below'} the target of {_pfmt(cur['budgetQty'])} units"
    elif has_prev:
        s += f", {'above' if cur['qty'] >= prev['qty'] else 'below'} last year's {_pfmt(prev['qty'])} units"
    return s


def _pnarrative_trend(nar, franchise):
    """Notes for the Monthly Sales trend slide — mirrors trendNarrativeHtml in
    sales_dashboard.js sentence-for-sentence (plain text, no bold/citations)."""
    if not nar:
        return "Not enough data to generate a summary for this selection."
    f_label = _pfranchise_label(franchise)
    month_label = f"{nar['monthName']} {nar['year']}"

    lines = []
    if nar.get('isPartialMonth') and nar.get('asOfDate'):
        as_of = fields.Date.from_string(nar['asOfDate'][:10])
        lines.append(
            f"Data for {month_label} is only synced through {as_of.strftime('%B %-d')} — the sales figures below "
            "are month-to-date, compared against the same number of days last year, while the budget is the "
            "full-month target."
        )

    lines.append(_pmonth_value_sentence(month_label, f_label, nar['month'], nar.get('monthPrevYear')) + ".")
    lines.append(_pmonth_qty_sentence(month_label, nar['month'], nar.get('monthPrevYear')) + ".")
    lines.append(f"Year-to-date {nar['year']} performance")
    lines.append(_pytd_value_sentence(nar['year'], f_label, nar['ytd'], nar.get('ytdPrevYear')) + ".")
    lines.append(_pytd_qty_sentence(nar['year'], nar['ytd'], nar.get('ytdPrevYear')) + ".")
    return "\n".join(lines)


def _pnarrative_category_single(data, dim_label, key='sales'):
    """Mirrors categoryNarrativeSingle in the on-screen dashboard's JS —
    used for region/franchise/customerType (not product or salesman, which
    have their own wording — see _pnarrative_product_* / _pnarrative_salesman_*)."""
    if not data:
        return f"No data available for {dim_label}."
    top = max(data, key=lambda d: d.get(key) or 0)
    total = sum((d.get(key) or 0) for d in data)
    text = (f"{top['label']} leads {dim_label} with sales of {_pfmt_m(top.get(key, 0))}, "
            f"accounting for {_pshare(top.get(key, 0), total)} of the total shown")
    if len(data) > 1:
        low = min(data, key=lambda d: d.get(key) or 0)
        text += f", while {low['label']} trails at {_pfmt_m(low.get(key, 0))}"
    return text + "."


def _pnarrative_category_single3(data, dim_label, is_partial_month, key='sales'):
    """Target / This Year / Last Year variant of _pnarrative_category_single,
    used once a single month has a prior-year figure to compare against —
    mirrors categoryNarrativeSingle3 in the on-screen dashboard's JS."""
    if not data:
        return f"No data available for {dim_label}."
    top = max(data, key=lambda d: d.get(key) or 0)
    total = sum((d.get(key) or 0) for d in data)
    achv = (top[key] / top['budget'] * 100) if top.get('budget') else None
    yoy = ((top.get(key, 0) - top['prevYearSales']) / top['prevYearSales'] * 100) if top.get('prevYearSales') else None
    target_word = 'full-month target' if is_partial_month else 'target'
    last_year_phrase = (f"the same days last year ({_pfmt_m(top.get('prevYearSales'))})" if is_partial_month
                         else f"last year's {_pfmt_m(top.get('prevYearSales'))}")
    text = (f"{top['label']} leads {dim_label} with sales of {_pfmt_m(top.get(key, 0))}, "
            f"accounting for {_pshare(top.get(key, 0), total)} of the total shown")
    if achv is not None:
        text += f", {'ahead of' if achv >= 100 else 'behind'} its {target_word} of {_pfmt_m(top.get('budget'))}"
    if yoy is not None:
        text += f", and {'up' if yoy >= 0 else 'down'} {abs(yoy):.0f}% versus {last_year_phrase}"
    if len(data) > 1:
        low = min(data, key=lambda d: d.get(key) or 0)
        text += f", while {low['label']} trails at {_pfmt_m(low.get(key, 0))}"
    return text + "."


def _pnarrative_category_compare(data, dim_label):
    if not data:
        return f"No data available for {dim_label}."
    top = max(data, key=lambda d: d.get('y2025') or 0)
    text = f"{top['label']} leads {dim_label} in 2025 with sales of {_pfmt_m(top.get('y2025', 0))}"
    growers = [d for d in data if d.get('y2024')]
    if growers:
        g = max(growers, key=lambda d: (d['y2025'] - d['y2024']) / d['y2024'])
        yoy = (g['y2025'] - g['y2024']) / g['y2024'] * 100
        text += f", and {g['label']} grew the most year-over-year at {yoy:+.0f}%"
    return text + "."


def _pnarrative_product_single(data):
    """Mirrors productNarrativeSingle — product groups get their own wording
    on screen (\"is the top-selling product group\"), distinct from the
    generic category phrasing above; data is expected pre-sorted (data[0])."""
    if not data:
        return "No data available for product groups."
    total = sum((d.get('value') or 0) for d in data)
    top = data[0]
    text = (f"{top['label']} is the top-selling product group at {_pfmt_m(top.get('value', 0))}, "
            f"representing {_pshare(top.get('value', 0), total)} of sales across the top {len(data)} groups shown")
    return text + "."


def _pnarrative_product_single3(data, is_partial_month):
    """Mirrors productNarrativeSingle3."""
    if not data:
        return "No data available for product groups."
    total = sum((d.get('value') or 0) for d in data)
    top = data[0]
    achv = (top['value'] / top['budget'] * 100) if top.get('budget') else None
    yoy = ((top.get('value', 0) - top['prevYearSales']) / top['prevYearSales'] * 100) if top.get('prevYearSales') else None
    target_word = 'full-month target' if is_partial_month else 'target'
    last_year_phrase = (f"the same days last year ({_pfmt_m(top.get('prevYearSales'))})" if is_partial_month
                         else f"last year's {_pfmt_m(top.get('prevYearSales'))}")
    text = (f"{top['label']} is the top-selling product group at {_pfmt_m(top.get('value', 0))}, "
            f"representing {_pshare(top.get('value', 0), total)} of sales across the top {len(data)} groups shown")
    if achv is not None:
        text += f", {'ahead of' if achv >= 100 else 'behind'} its {target_word} of {_pfmt_m(top.get('budget'))}"
    if yoy is not None:
        text += f", and {'up' if yoy >= 0 else 'down'} {abs(yoy):.0f}% versus {last_year_phrase}"
    return text + "."


def _pnarrative_product_compare(data):
    """Mirrors productNarrativeCompare — takes data[0] (not a max()), matching
    the JS, since merge_dim's pre-sort already orders by combined volume."""
    if not data:
        return "No data available for product groups."
    top = data[0]
    yoy = ((top['y2025'] - top['y2024']) / top['y2024'] * 100) if top.get('y2024') else None
    yoy_phrase = ("with no comparable 2024 figure" if yoy is None
                  else f"{'up' if yoy >= 0 else 'down'} {abs(yoy):.0f}% year-over-year")
    return f"{top['label']} remains the top-selling product group in 2025 at {_pfmt_m(top.get('y2025', 0))}, {yoy_phrase}."


def _pnarrative_salesman_single(data):
    """Mirrors salesmanNarrativeSingle — salesmen get their own wording
    (\"is the top performer\"), distinct from the category phrasing above."""
    if not data:
        return "No data available for salesmen."
    top = max(data, key=lambda d: d.get('sales') or 0)
    achvs = [(d['sales'] / d['budget'] * 100) for d in data if d.get('budget')]
    avg = sum(achvs) / len(achvs) if achvs else None
    text = f"{top['label']} is the top performer with sales of {_pfmt_m(top.get('sales', 0))}"
    if avg is not None:
        text += f", and the group is averaging {avg:.0f}% achievement against budget"
    return text + "."


def _pnarrative_salesman_single3(data, is_partial_month):
    """Mirrors salesmanNarrativeSingle3."""
    if not data:
        return "No data available for salesmen."
    top = max(data, key=lambda d: d.get('sales') or 0)
    achvs = [(d['sales'] / d['budget'] * 100) for d in data if d.get('budget')]
    avg_achv = sum(achvs) / len(achvs) if achvs else None
    yoys = [((d['sales'] - d['prevYearSales']) / d['prevYearSales'] * 100) for d in data if d.get('prevYearSales')]
    avg_yoy = sum(yoys) / len(yoys) if yoys else None
    target_word = 'full-month target' if is_partial_month else 'target'
    last_year_word = 'the same days last year' if is_partial_month else 'last year'
    text = f"{top['label']} is the top performer with sales of {_pfmt_m(top.get('sales', 0))}"
    if avg_achv is not None:
        text += f", and the group is averaging {avg_achv:.0f}% achievement against {target_word}"
    if avg_yoy is not None:
        text += f", {'up' if avg_yoy >= 0 else 'down'} {abs(avg_yoy):.0f}% versus {last_year_word} on average"
    return text + "."


def _pnarrative_salesman_compare(data):
    """Mirrors salesmanNarrativeCompare (signed +/-% average, not word form)."""
    if not data:
        return "No data available for salesmen."
    top = max(data, key=lambda d: d.get('y2025') or 0)
    yoys = [((d['y2025'] - d['y2024']) / d['y2024'] * 100) for d in data if d.get('y2024')]
    avg = sum(yoys) / len(yoys) if yoys else None
    text = f"{top['label']} leads in 2025 with sales of {_pfmt_m(top.get('y2025', 0))}"
    if avg is not None:
        text += f", with the group averaging {avg:+.0f}% year-over-year"
    return text + "."


class PbiSalesDashboardController(http.Controller):

    # ------------------------------------------------------------------
    # low-level helpers
    # ------------------------------------------------------------------
    def _query(self, sql, params=()):
        request.env.cr.execute(sql, params)
        cols = [d[0] for d in request.env.cr.description]
        return [dict(zip(cols, row)) for row in request.env.cr.fetchall()]

    def _parse_year_param(self, year):
        """Returns (mode, year_int, month_int). mode is one of
        'all', 'compare', 'year' (year_int only) or 'month' (year_int + month_int)."""
        if year in ("all", "compare"):
            return year, None, None
        m = MONTH_VALUE_RE.match(str(year))
        if m:
            return "month", int(m.group(1)), int(m.group(2))
        return "year", int(year), None

    def _dim_clauses(self, franchise, customer_type, region):
        """Clause fragments + params for the three cross-cutting dimension
        filters (Franchise, Customer Type, Region) — shared by every query
        builder below so a selection narrows the whole dashboard, exactly
        like Franchise already did before Customer Type/Region existed."""
        clauses, params = [], []
        if franchise and franchise != "all":
            clauses.append("bi_franchisename = %s")
            params.append(franchise)
        if customer_type and customer_type != "all":
            clauses.append("bi_csttypedesc = %s")
            params.append(customer_type)
        if region and region != "all":
            clauses.append("bi_cstregiondesc = %s")
            params.append(region)
        return clauses, params

    def _where(self, year, franchise, customer_type='all', region='all'):
        clauses, params = [], []
        mode, y, m = self._parse_year_param(year)
        if mode in ("year", "month"):
            clauses.append("bi_year = %s")
            params.append(y)
        if mode == "month":
            clauses.append("bi_month = %s")
            params.append(m)
        dim_clauses, dim_params = self._dim_clauses(franchise, customer_type, region)
        clauses += dim_clauses
        params += dim_params
        return (" AND " + " AND ".join(clauses)) if clauses else "", params

    def _where_year_only(self, year, franchise, customer_type='all', region='all'):
        """Like _where but ignores any month restriction — used for the
        monthly trend chart, which should keep showing the full year even
        when a single month is selected, for context."""
        clauses, params = [], []
        mode, y, m = self._parse_year_param(year)
        if mode in ("year", "month"):
            clauses.append("bi_year = %s")
            params.append(y)
        dim_clauses, dim_params = self._dim_clauses(franchise, customer_type, region)
        clauses += dim_clauses
        params += dim_params
        return (" AND " + " AND ".join(clauses)) if clauses else "", params

    def _latest_data_month(self, franchise, customer_type='all', region='all'):
        """(year, month) of the most recent month with any actual sales row,
        for the given filters. Used to tell a month that is still syncing in
        (today's month) apart from an ordinary completed month whose
        calendar-last-day simply had no transactions (weekend/holiday) —
        the latter must NOT be treated as partial."""
        dim_clauses, dim_params = self._dim_clauses(franchise, customer_type, region)
        fclause = (" AND " + " AND ".join(dim_clauses)) if dim_clauses else ""
        row = self._query(f"""SELECT bi_year, bi_month FROM v_bidata_live
                              WHERE bi_amount IS NOT NULL {fclause}
                              ORDER BY bi_year DESC, bi_month DESC LIMIT 1""", dim_params)
        return (row[0]["bi_year"], row[0]["bi_month"]) if row else None

    def _month_sync_info(self, franchise, year, month, customer_type='all', region='all'):
        """(as_of_date, day_lte) for (year, month). day_lte is only set when
        this is the genuinely latest-syncing month for these filters AND
        it isn't complete yet — callers then compare against the same day
        range last year. For any other month (including a historical month
        whose calendar-last-day simply had no transactions) both are None
        so callers compare full months as before."""
        if self._latest_data_month(franchise, customer_type, region) != (year, month):
            return None, None
        dim_clauses, dim_params = self._dim_clauses(franchise, customer_type, region)
        fclause = (" AND " + " AND ".join(dim_clauses)) if dim_clauses else ""
        row = self._query(f"""SELECT max(bi_monthdate) AS d FROM v_bidata_live
                              WHERE bi_year = %s AND bi_month = %s AND bi_amount IS NOT NULL {fclause}""",
                           [year, month] + dim_params)
        as_of = row[0]["d"]
        if not as_of:
            return None, None
        last_day = calendar.monthrange(year, month)[1]
        day_lte = as_of.day if as_of.day < last_day else None
        return as_of, day_lte

    def _prev_year_dim_sales(self, dim_expr, franchise, year, month, day_lte, customer_type='all', region='all',
                              extra_where="", extra_params=None):
        """{label: sales} for a breakdown dimension in a single prior-year
        month, day-limited like _period_sums when day_lte is given. Used to
        attach a 'Last Year' figure onto this year's region/franchise/
        customer-type/product/salesman/drill-down rows."""
        clauses, params = ["bi_year = %s", "bi_month = %s"], [year, month]
        if day_lte is not None:
            clauses.append("extract(day from bi_monthdate) <= %s")
            params.append(day_lte)
        dim_clauses, dim_params = self._dim_clauses(franchise, customer_type, region)
        clauses += dim_clauses
        params += dim_params
        where = " AND ".join(clauses) + extra_where
        if extra_params:
            params = params + extra_params
        rows = self._query(f"""
            SELECT {dim_expr} AS label, round(sum(bi_amount)::numeric,0) AS sales
            FROM v_bidata_live WHERE {dim_expr} IS NOT NULL AND {where}
            GROUP BY 1
        """, params)
        return {r["label"]: (float(r["sales"]) if r["sales"] is not None else 0.0) for r in rows}

    def _period_sums(self, year, franchise, month=None, month_lte=None, day_lte=None, customer_type='all', region='all'):
        clauses, params = ["bi_year = %s"], [year]
        if month is not None:
            clauses.append("bi_month = %s")
            params.append(month)
        if month_lte is not None:
            clauses.append("bi_month <= %s")
            params.append(month_lte)
        if day_lte is not None:
            # restrict to the same day-of-month range, so a partial (still
            # syncing) current month is compared like-for-like against the
            # same number of days last year, instead of a full prior month.
            clauses.append("extract(day from bi_monthdate) <= %s")
            params.append(day_lte)
        dim_clauses, dim_params = self._dim_clauses(franchise, customer_type, region)
        clauses += dim_clauses
        params += dim_params
        where = " AND ".join(clauses)
        r = self._query(f"""SELECT sum(bi_amount) AS sales, sum(bi_budgetamount) AS budget,
                                   sum(bi_qty) AS qty, sum(bi_budgetqty) AS budgetqty
                            FROM v_bidata_live WHERE {where}""", params)[0]
        return {"sales": float(r["sales"] or 0), "budget": float(r["budget"] or 0),
                "qty": int(r["qty"] or 0), "budgetQty": int(r["budgetqty"] or 0)}

    # ------------------------------------------------------------------
    # narrative (latest month + YTD vs target vs prior year)
    # ------------------------------------------------------------------
    def _fetch_narrative(self, year, franchise, customer_type='all', region='all'):
        mode, y, m = self._parse_year_param(year)

        if mode == "month":
            ref_year, latest_month = y, m
        else:
            ref_year = y if mode == "year" else 2025  # "all" falls back to the latest year

            dim_clauses, dim_params = self._dim_clauses(franchise, customer_type, region)
            fclause = (" AND " + " AND ".join(dim_clauses)) if dim_clauses else ""
            row = self._query(f"""SELECT max(bi_month) AS m FROM v_bidata_live
                                   WHERE bi_year = %s AND bi_amount IS NOT NULL {fclause}""",
                               [ref_year] + dim_params)
            latest_month = row[0]["m"]
            if latest_month is None:
                return None

        has_prev_year = (ref_year - 1) >= 2023
        as_of, day_lte = self._month_sync_info(franchise, ref_year, latest_month, customer_type, region)

        return {
            "year": ref_year,
            "monthNum": latest_month,
            "monthName": MONTH_NAMES[latest_month - 1],
            "month": self._period_sums(ref_year, franchise, month=latest_month, customer_type=customer_type, region=region),
            "monthPrevYear": self._period_sums(ref_year - 1, franchise, month=latest_month, day_lte=day_lte,
                                                customer_type=customer_type, region=region) if has_prev_year else None,
            "ytd": self._period_sums(ref_year, franchise, month_lte=latest_month, customer_type=customer_type, region=region),
            "ytdPrevYear": self._period_sums(ref_year - 1, franchise, month_lte=latest_month,
                                              customer_type=customer_type, region=region) if has_prev_year else None,
            "isPartialMonth": day_lte is not None,
            "asOfDate": as_of.isoformat() if as_of else None,
        }

    # ------------------------------------------------------------------
    # single-year / all-years bundle
    # ------------------------------------------------------------------
    def _fetch_bundle(self, year, franchise, customer_type='all', region='all'):
        mode, y, m = self._parse_year_param(year)
        where, params = self._where(year, franchise, customer_type, region)

        kpi_rows = self._query(f"""
            SELECT round(sum(bi_amount)::numeric,0) AS sales,
                   round(sum(bi_budgetamount)::numeric,0) AS budget,
                   sum(bi_qty) AS qty, sum(bi_budgetqty) AS budget_qty,
                   count(distinct bi_cstno) AS customers
            FROM v_bidata_live WHERE 1=1 {where}
        """, params)
        k = kpi_rows[0]
        kpis = {"sales": float(k["sales"] or 0), "budget": float(k["budget"] or 0),
                "qty": int(k["qty"] or 0), "budgetQty": int(k["budget_qty"] or 0),
                "customers": int(k["customers"] or 0)}

        # Same-period-last-year sales, for the single-month view only. Only
        # the single genuinely-latest-syncing month gets the day-limited,
        # like-for-like comparison — an older completed month whose
        # calendar-last-day simply had no transactions (weekend/holiday)
        # compares as a full month, same as always.
        day_lte = None
        has_prev_year_dims = False
        if mode == "month":
            as_of, day_lte = self._month_sync_info(franchise, y, m, customer_type, region)
            if as_of:
                kpis["asOfDate"] = as_of.isoformat()
                kpis["isPartialMonth"] = day_lte is not None

            if (y - 1) >= 2023:
                prev = self._period_sums(y - 1, franchise, month=m, day_lte=day_lte,
                                          customer_type=customer_type, region=region)
                kpis["prevYearSales"] = prev["sales"]
                has_prev_year_dims = True

        region_rows = self._query(f"""
            SELECT bi_cstregiondesc AS label, round(sum(bi_amount)::numeric,0) AS sales,
                   round(sum(bi_budgetamount)::numeric,0) AS budget
            FROM v_bidata_live WHERE bi_cstregiondesc IS NOT NULL {where}
            GROUP BY 1 ORDER BY 2 DESC
        """, params)

        franchise_bd = self._query(f"""
            SELECT bi_franchisename AS label, round(sum(bi_amount)::numeric,0) AS sales,
                   round(sum(bi_budgetamount)::numeric,0) AS budget
            FROM v_bidata_live WHERE bi_franchisename IS NOT NULL {where}
            GROUP BY 1 HAVING sum(bi_amount) > 0 ORDER BY 2 DESC LIMIT 10
        """, params)

        customer_type_rows = self._query(f"""
            SELECT bi_csttypedesc AS label, round(sum(bi_amount)::numeric,0) AS sales,
                   round(sum(bi_budgetamount)::numeric,0) AS budget
            FROM v_bidata_live WHERE bi_csttypedesc IS NOT NULL {where}
            GROUP BY 1 HAVING sum(bi_amount) > 0 ORDER BY 2 DESC LIMIT 10
        """, params)

        product = self._query(f"""
            SELECT bi_pgroupname AS label, round(sum(bi_amount)::numeric,0) AS value,
                   round(sum(bi_budgetamount)::numeric,0) AS budget
            FROM v_bidata_live WHERE bi_pgroupname IS NOT NULL {where}
            GROUP BY 1 HAVING sum(bi_amount) > 0 ORDER BY 2 DESC LIMIT 10
        """, params)

        salesman = self._query(f"""
            SELECT bi_salesmanname AS label, round(sum(bi_amount)::numeric,0) AS sales,
                   round(sum(bi_budgetamount)::numeric,0) AS budget
            FROM v_bidata_live WHERE bi_salesmanname IS NOT NULL {where}
            GROUP BY 1 HAVING sum(bi_amount) > 0 ORDER BY 2 DESC LIMIT 10
        """, params)

        trend_where, trend_params = self._where_year_only(year, franchise, customer_type, region)
        monthly = self._query(f"""
            SELECT to_char(bi_monthdate,'YYYY-MM') AS label,
                   round(sum(bi_amount)::numeric,0) AS sales,
                   round(sum(bi_budgetamount)::numeric,0) AS budget
            FROM v_bidata_live WHERE bi_monthdate IS NOT NULL {trend_where}
            GROUP BY 1 ORDER BY 1
        """, trend_params)

        def num(v):
            return None if v is None else float(v)

        region_list = [{"label": r["label"], "sales": num(r["sales"]), "budget": num(r["budget"])} for r in region_rows]
        franchise_list = [{"label": r["label"], "sales": num(r["sales"]), "budget": num(r["budget"])} for r in franchise_bd]
        customer_type_list = [{"label": r["label"], "sales": num(r["sales"]), "budget": num(r["budget"])} for r in customer_type_rows]
        product_list = [{"label": r["label"], "value": num(r["value"]), "budget": num(r["budget"])} for r in product]
        salesman_list = [{"label": r["label"], "sales": num(r["sales"]), "budget": num(r["budget"])} for r in salesman]
        monthly_list = [{"label": r["label"], "sales": num(r["sales"]), "budget": num(r["budget"])} for r in monthly]

        if has_prev_year_dims:
            for rows, dim_expr in (
                (region_list, "bi_cstregiondesc"), (franchise_list, "bi_franchisename"),
                (customer_type_list, "bi_csttypedesc"), (product_list, "bi_pgroupname"),
                (salesman_list, "bi_salesmanname"),
            ):
                prev_map = self._prev_year_dim_sales(dim_expr, franchise, y - 1, m, day_lte, customer_type, region)
                for r in rows:
                    r["prevYearSales"] = prev_map.get(r["label"], 0.0)

            # Same treatment for the trend line: the whole prior year's
            # monthly sales, so the chart can draw a full "Last Year" line
            # alongside this year's Sales/Budget, not just the one dot.
            prev_trend_where, prev_trend_params = self._where_year_only(str(y - 1), franchise, customer_type, region)
            prev_monthly = self._query(f"""
                SELECT to_char(bi_monthdate,'YYYY-MM') AS label, round(sum(bi_amount)::numeric,0) AS sales
                FROM v_bidata_live WHERE bi_monthdate IS NOT NULL {prev_trend_where}
                GROUP BY 1
            """, prev_trend_params)
            prev_by_month = {int(r["label"].split("-")[1]): num(r["sales"]) for r in prev_monthly}
            for row in monthly_list:
                row["prevYearSales"] = prev_by_month.get(int(row["label"].split("-")[1]), 0.0)

        return {
            "kpis": kpis,
            "region": region_list,
            "franchise": franchise_list,
            "customerType": customer_type_list,
            "product": product_list,
            "salesman": salesman_list,
            "monthly": monthly_list,
            "narrative": self._fetch_narrative(year, franchise, customer_type, region),
        }

    # ------------------------------------------------------------------
    # 2024 vs 2025 compare bundle
    # ------------------------------------------------------------------
    def _fetch_compare(self, franchise, customer_type='all', region='all'):
        y2024 = self._fetch_bundle("2024", franchise, customer_type, region)
        y2025 = self._fetch_bundle("2025", franchise, customer_type, region)

        def by_month_number(monthly, key):
            out = {}
            for m in monthly:
                mm = int(m["label"].split("-")[1])
                out[mm] = m[key]
            return out

        m24 = by_month_number(y2024["monthly"], "sales")
        m25 = by_month_number(y2025["monthly"], "sales")
        b25 = by_month_number(y2025["monthly"], "budget")
        months = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
        trend = [{"label": months[i - 1], "y2024": m24.get(i), "y2025": m25.get(i), "budget2025": b25.get(i)}
                 for i in range(1, 13)]

        def merge_dim(rows24, rows25, key, budget_key=None):
            labels = []
            for r in rows24 + rows25:
                if r["label"] not in labels:
                    labels.append(r["label"])
            v24 = {r["label"]: r[key] for r in rows24}
            v25 = {r["label"]: r[key] for r in rows25}
            bud25 = {r["label"]: r[budget_key] for r in rows25} if budget_key else {}
            merged = [{"label": l, "y2024": v24.get(l, 0) or 0, "y2025": v25.get(l, 0) or 0,
                       "budget2025": bud25.get(l, 0) or 0} for l in labels]
            merged.sort(key=lambda d: (d["y2024"] or 0) + (d["y2025"] or 0), reverse=True)
            return merged[:10]

        return {
            "kpis2024": y2024["kpis"], "kpis2025": y2025["kpis"],
            "trend": trend,
            "region": merge_dim(y2024["region"], y2025["region"], "sales", "budget"),
            "franchise": merge_dim(y2024["franchise"], y2025["franchise"], "sales", "budget"),
            "customerType": merge_dim(y2024["customerType"], y2025["customerType"], "sales", "budget"),
            "product": merge_dim(y2024["product"], y2025["product"], "value", "budget"),
            "salesman": merge_dim(y2024["salesman"], y2025["salesman"], "sales", "budget"),
            "narrative": y2025["narrative"],
        }

    # ------------------------------------------------------------------
    # PowerPoint export — native editable charts + editable speaker notes
    # ------------------------------------------------------------------
    def _pptx_fix_notes_master_ref(self, prs):
        """python-pptx creates the notesMaster part + relationship the first
        time any slide.notes_slide is accessed, but never adds the
        <p:notesMasterIdLst> pointer to presentation.xml itself (a known
        python-pptx gap). Without it, the package is inconsistent — some
        viewers fail to show the Notes pane at all, and PowerPoint can flag
        the file as needing repair, which can also blank/garble other shapes
        (e.g. slide titles) during that recovery pass. Add the missing
        pointer by hand if a notesMaster relationship exists.
        """
        presentation_part = prs.part
        notes_master_rid = None
        for rid, rel in presentation_part.rels.items():
            if rel.reltype.endswith("/notesMaster"):
                notes_master_rid = rid
                break
        if not notes_master_rid:
            return  # no notes slide was ever created; nothing to fix

        p_elm = prs.element
        if p_elm.find(qn("p:notesMasterIdLst")) is not None:
            return  # already present

        notes_master_id_lst = p_elm.makeelement(qn("p:notesMasterIdLst"), {})
        notes_master_id = notes_master_id_lst.makeelement(
            qn("p:notesMasterId"), {qn("r:id"): notes_master_rid}
        )
        notes_master_id_lst.append(notes_master_id)
        # Schema order requires notesMasterIdLst right after sldMasterIdLst.
        sld_master_id_lst = p_elm.find(qn("p:sldMasterIdLst"))
        sld_master_id_lst.addnext(notes_master_id_lst)

    def _pptx_add_title_slide(self, prs, layout, year, franchise, customer_type, region, kpis_text):
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(12.1), Inches(1.2))
        p = tb.text_frame.paragraphs[0]
        p.text = "BI Data Live — Sales Dashboard"
        p.font.size = Pt(36)
        p.font.bold = True

        franchise_label = "All Franchises" if franchise == "all" else franchise
        customer_type_label = "All Customer Types" if customer_type == "all" else customer_type
        region_label = "All Regions" if region == "all" else region
        month_match = MONTH_VALUE_RE.match(year)
        if year == "compare":
            year_label = "Compare 2024 vs 2025"
        elif year == "all":
            year_label = "All Years (2023–2025)"
        elif month_match:
            year_label = f"{MONTH_NAMES[int(month_match.group(2)) - 1]} {month_match.group(1)}"
        else:
            year_label = f"Year {year}"

        sub = slide.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(12.1), Inches(1.2))
        sp = sub.text_frame.paragraphs[0]
        sp.text = f"Source: v_bidata_live  ·  {franchise_label}  ·  {customer_type_label}  ·  {region_label}  ·  {year_label}"
        sp.font.size = Pt(16)
        gen = sub.text_frame.add_paragraph()
        gen.text = f"Generated {fields.Datetime.now().strftime('%Y-%m-%d %H:%M')}"
        gen.font.size = Pt(12)

        kpi_box = slide.shapes.add_textbox(Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.2))
        kpi_box.text_frame.word_wrap = True
        kp = kpi_box.text_frame.paragraphs[0]
        kp.text = kpis_text
        kp.font.size = Pt(15)
        return slide

    def _pptx_style_series(self, chart, colors):
        for i, series in enumerate(chart.plots[0].series):
            if i < len(colors):
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = colors[i]

    def _pptx_add_chart_slide(self, prs, layout, title, notes, chart_type, categories, series, colors, subtitle=None):
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        if subtitle:
            st = slide.shapes.add_textbox(Inches(0.4), Inches(0.7), Inches(12.5), Inches(0.4))
            sp = st.text_frame.paragraphs[0]
            sp.text = subtitle
            sp.font.size = Pt(12)
            sp.font.color.rgb = RGBColor(0x89, 0x87, 0x81)

        chart_data = CategoryChartData()
        chart_data.categories = categories
        for name, values in series:
            chart_data.add_series(name, values)
        gframe = slide.shapes.add_chart(chart_type, Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.8), chart_data)
        chart = gframe.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        self._pptx_style_series(chart, colors)

        if notes:
            _pptx_write_notes(slide, notes)
        return slide

    def _pptx_add_table_slide(self, prs, layout, title, notes, headers, rows):
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True

        n_rows, n_cols = len(rows) + 1, len(headers)
        table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.4), Inches(1.1), Inches(12.5), Inches(5.8))
        table = table_shape.table
        for c, h in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = h
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
        for r_i, row in enumerate(rows, start=1):
            for c_i, val in enumerate(row):
                cell = table.cell(r_i, c_i)
                cell.text = str(val)
                cell.text_frame.paragraphs[0].font.size = Pt(11)

        if notes:
            _pptx_write_notes(slide, notes)
        return slide

    def _build_pptx(self, year, franchise, customer_type='all', region='all'):
        compare = year == "compare"
        bundle = (self._fetch_compare(franchise, customer_type, region) if compare
                  else self._fetch_bundle(year, franchise, customer_type, region))

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        if compare:
            self._pptx_add_title_slide(
                prs, blank, year, franchise, customer_type, region,
                _pnarrative_kpis_compare(bundle["kpis2024"], bundle["kpis2025"]),
            )

            cats = [d["label"] for d in bundle["trend"]]
            series = [
                ("2025 Budget", [d.get("budget2025") or 0 for d in bundle["trend"]]),
                ("2025 Sales", [d.get("y2025") or 0 for d in bundle["trend"]]),
                ("2024 Sales", [d.get("y2024") or 0 for d in bundle["trend"]]),
            ]
            self._pptx_add_chart_slide(
                prs, blank, "Monthly Sales — 2024 vs 2025",
                _pnarrative_trend(bundle["narrative"], franchise),
                XL_CHART_TYPE.LINE_MARKERS, cats, series,
                [PPTX_COLOR_BUDGET, PPTX_COLOR_2025, PPTX_COLOR_2024],
            )

            def dim_slide(key, title, dim_label):
                data = bundle[key]
                cats = [d["label"] for d in data]
                series = [
                    ("2025 Budget", [d.get("budget2025") or 0 for d in data]),
                    ("2025 Sales", [d.get("y2025") or 0 for d in data]),
                    ("2024 Sales", [d.get("y2024") or 0 for d in data]),
                ]
                is_product = key == "product"
                notes = _pnarrative_product_compare(data) if is_product else _pnarrative_category_compare(data, dim_label)
                chart_type = XL_CHART_TYPE.BAR_CLUSTERED if is_product else XL_CHART_TYPE.COLUMN_CLUSTERED
                self._pptx_add_chart_slide(
                    prs, blank, f"{title} — 2024 vs 2025", notes,
                    chart_type, cats, series, [PPTX_COLOR_BUDGET, PPTX_COLOR_2025, PPTX_COLOR_2024],
                )

            dim_slide("region", "Sales by Region", "regions")
            dim_slide("franchise", "Sales by Franchise", "franchises")
            dim_slide("customerType", "Sales by Customer Type", "customer types")
            dim_slide("product", "Top Product Groups", "product groups")

            sm = bundle["salesman"]
            rows = [[d["label"], _pfmt(d.get("budget2025")), _pfmt(d.get("y2025")), _pfmt(d.get("y2024"))] for d in sm]
            self._pptx_add_table_slide(
                prs, blank, "Top 10 Salesmen — 2024 vs 2025",
                _pnarrative_salesman_compare(sm),
                ["Salesman", "Budget 2025", "Sales 2025", "Sales 2024"], rows,
            )
        else:
            kpis = bundle["kpis"]
            # Only present once a single month is selected AND its prior-year
            # figure exists — mirrors hasPrevYear/isPartialMonth in the
            # on-screen JS (see setMode/load in sales_dashboard.js), so the
            # exported deck always matches what's on screen for that period.
            has_prev_year = kpis.get("prevYearSales") is not None
            is_partial_month = bool(kpis.get("isPartialMonth"))
            month_match = MONTH_VALUE_RE.match(year)
            suffix = f" — {MONTH_NAMES[int(month_match.group(2)) - 1][:3]} {month_match.group(1)}" if month_match else ""

            self._pptx_add_title_slide(prs, blank, year, franchise, customer_type, region, _pnarrative_kpis_single(kpis))

            target_label = "Target (Full Month)" if is_partial_month else "Target"
            this_year_label = "This Year (MTD)" if is_partial_month else "This Year"
            last_year_label = "Last Year (Same Days)" if is_partial_month else "Last Year"

            cats = [d["label"] for d in bundle["monthly"]]
            if has_prev_year:
                trend_series = [
                    (this_year_label, [d.get("sales") or 0 for d in bundle["monthly"]]),
                    (target_label, [d.get("budget") or 0 for d in bundle["monthly"]]),
                    (last_year_label, [d.get("prevYearSales") or 0 for d in bundle["monthly"]]),
                ]
                trend_colors = [PPTX_COLOR_2025, PPTX_COLOR_BUDGET, PPTX_COLOR_2024]
            else:
                trend_series = [
                    ("Sales", [d.get("sales") or 0 for d in bundle["monthly"]]),
                    ("Budget", [d.get("budget") or 0 for d in bundle["monthly"]]),
                ]
                trend_colors = [PPTX_COLOR_2025, PPTX_COLOR_BUDGET]
            self._pptx_add_chart_slide(
                prs, blank, f"Monthly Sales vs Budget{suffix}", _pnarrative_trend(bundle["narrative"], franchise),
                XL_CHART_TYPE.LINE_MARKERS, cats, trend_series, trend_colors,
            )

            def dim_slide(key, title, dim_label, value_key="sales"):
                data = bundle[key]
                cats = [d["label"] for d in data]
                is_product = key == "product"
                if has_prev_year:
                    series = [
                        (target_label, [d.get("budget") or 0 for d in data]),
                        (this_year_label, [d.get(value_key) or 0 for d in data]),
                        (last_year_label, [d.get("prevYearSales") or 0 for d in data]),
                    ]
                    colors = [PPTX_COLOR_BUDGET, PPTX_COLOR_2025, PPTX_COLOR_2024]
                    notes = (_pnarrative_product_single3(data, is_partial_month) if is_product
                             else _pnarrative_category_single3(data, dim_label, is_partial_month, value_key))
                else:
                    series = [
                        ("Sales", [d.get(value_key) or 0 for d in data]),
                        ("Budget", [d.get("budget") or 0 for d in data]),
                    ]
                    colors = [PPTX_COLOR_2025, PPTX_COLOR_BUDGET]
                    notes = (_pnarrative_product_single(data) if is_product
                             else _pnarrative_category_single(data, dim_label, value_key))
                chart_type = XL_CHART_TYPE.BAR_CLUSTERED if is_product else XL_CHART_TYPE.COLUMN_CLUSTERED
                self._pptx_add_chart_slide(
                    prs, blank, f"{title}{suffix}", notes,
                    chart_type, cats, series, colors,
                )

            dim_slide("region", "Sales by Region", "regions")
            dim_slide("franchise", "Sales by Franchise", "franchises")
            dim_slide("customerType", "Sales by Customer Type", "customer types")
            dim_slide("product", "Top Product Groups", "product groups", value_key="value")

            sm = bundle["salesman"]
            if has_prev_year:
                rows = []
                for d in sm:
                    prev = d.get("prevYearSales")
                    yoy = ((d.get("sales") or 0) - prev) / prev * 100 if prev else None
                    rows.append([d["label"], _pfmt(d.get("budget")), _pfmt(d.get("sales")), _pfmt(prev),
                                 f"{yoy:+.0f}%" if yoy is not None else "–"])
                self._pptx_add_table_slide(
                    prs, blank, f"Top 10 Salesmen{suffix}",
                    _pnarrative_salesman_single3(sm, is_partial_month),
                    ["Salesman", target_label, this_year_label, last_year_label, "YoY %"], rows,
                )
            else:
                rows = []
                for d in sm:
                    achv = (d["sales"] / d["budget"] * 100) if d.get("budget") else None
                    rows.append([d["label"], _pfmt(d.get("sales")), _pfmt(d.get("budget")), f"{achv:.0f}%" if achv is not None else "–"])
                self._pptx_add_table_slide(
                    prs, blank, f"Top 10 Salesmen{suffix}", _pnarrative_salesman_single(sm),
                    ["Salesman", "Sales Amount", "Budget Amount", "Achv %"], rows,
                )

        # Stashed in a document property (invisible in the UI, untouched by
        # editing slide content/notes) so import_notes knows which period and
        # franchise this file was actually exported for — it must NOT rely on
        # whatever happens to be selected in the browser at import time, since
        # the user may have switched screens between exporting and importing.
        self._pptx_fix_notes_master_ref(prs)
        prs.core_properties.keywords = (
            f"pbi_dashboards|year={year}|franchise={franchise}|customerType={customer_type}|region={region}")
        stream = BytesIO()
        prs.save(stream)
        return stream.getvalue()

    @http.route('/pbi_dashboards/sales/export.pptx', type='http', auth='user')
    def export_pptx(self, year='2025', franchise='Midea', customerType='all', region='all', **kwargs):
        if not self._has_access():
            return request.not_found()
        try:
            data = self._build_pptx(year, franchise, customerType, region)
        except Exception as e:
            return request.make_response(f"Error building PPTX: {e}", status=500)
        headers = [
            ('Content-Type', 'application/vnd.openxmlformats-officedocument.presentationml.presentation'),
            ('Content-Disposition', 'attachment; filename="BI_Data_Live_Sales_Dashboard.pptx"'),
            ('Content-Length', len(data)),
        ]
        return request.make_response(data, headers=headers)

    # ------------------------------------------------------------------
    # region drill-down: Region -> City -> Salesman
    # ------------------------------------------------------------------
    # v_bidata_live has no clean city column; bi_invwhousename ("Riyadh Main
    # WH", "Jeddah Showroom", ...) is the closest reliable proxy for a
    # customer/branch's city, so "City" here means that field with its
    # warehouse/showroom/damage suffix stripped off.
    CITY_EXPR = (
        r"regexp_replace(bi_invwhousename, "
        r"'\s*(Main WH|Damage Cartons|Damage|Showroom)\s*$', '', 'i')"
    )

    def _region_total_budget(self, year, franchise, customer_type, region, drill_region):
        """Real budget total for one drilled-into region, used to derive a
        per-city estimate (see _drill_dimension below) since budget rows have
        no warehouse/city attribution of their own."""
        where, params = self._where(year, franchise, customer_type, region)
        row = self._query(f"""
            SELECT round(sum(bi_budgetamount)::numeric,0) AS budget
            FROM v_bidata_live WHERE bi_cstregiondesc = %s {where}
        """, [drill_region] + params)[0]
        return float(row["budget"] or 0)

    def _drill_dimension(self, year, franchise, customer_type, region, level, drill_region=None, drill_city=None):
        where, params = self._where(year, franchise, customer_type, region)
        if level in ("city", "salesman") and drill_region:
            where += " AND bi_cstregiondesc = %s"
            params.append(drill_region)
        if level == "salesman" and drill_city:
            where += f" AND {self.CITY_EXPR} = %s"
            params.append(drill_city)

        dim_expr = {
            "region": "bi_cstregiondesc",
            "city": self.CITY_EXPR,
            "salesman": "bi_salesmanname",
        }[level]

        rows = self._query(f"""
            SELECT {dim_expr} AS label, round(sum(bi_amount)::numeric,0) AS sales,
                   round(sum(bi_budgetamount)::numeric,0) AS budget
            FROM v_bidata_live WHERE {dim_expr} IS NOT NULL {where}
            GROUP BY 1 HAVING sum(bi_amount) > 0 ORDER BY 2 DESC LIMIT 15
        """, params)
        result = [{"label": r["label"], "sales": float(r["sales"] or 0),
                   "budget": float(r["budget"] or 0)} for r in rows]

        if level == "city":
            # Budget rows carry no warehouse name at all, so sum(bi_budgetamount)
            # above is always 0 per city. Estimate one instead by splitting the
            # parent region's real target across cities in proportion to each
            # city's share of sales, rather than leaving the target blank.
            region_budget = (self._region_total_budget(year, franchise, customer_type, region, drill_region)
                              if drill_region else 0.0)
            total_sales = sum(r["sales"] for r in result)
            for r in result:
                r["budget"] = round(region_budget * (r["sales"] / total_sales), 0) if total_sales else 0.0

        return result

    def _fetch_region_drill(self, year, franchise, customer_type, region, level, drill_region, drill_city):
        mode, y, m = self._parse_year_param(year)
        if mode != "compare":
            rows = self._drill_dimension(year, franchise, customer_type, region, level, drill_region, drill_city)
            if mode == "month" and (y - 1) >= 2023:
                _, day_lte = self._month_sync_info(franchise, y, m, customer_type, region)
                dim_expr = {"region": "bi_cstregiondesc", "city": self.CITY_EXPR, "salesman": "bi_salesmanname"}[level]
                extra_where, extra_params = "", []
                if level in ("city", "salesman") and drill_region:
                    extra_where += " AND bi_cstregiondesc = %s"
                    extra_params.append(drill_region)
                if level == "salesman" and drill_city:
                    extra_where += f" AND {self.CITY_EXPR} = %s"
                    extra_params.append(drill_city)
                prev_map = self._prev_year_dim_sales(dim_expr, franchise, y - 1, m, day_lte, customer_type, region,
                                                      extra_where, extra_params)
                for r in rows:
                    r["prevYearSales"] = prev_map.get(r["label"], 0.0)
            return rows

        rows24 = self._drill_dimension("2024", franchise, customer_type, region, level, drill_region, drill_city)
        rows25 = self._drill_dimension("2025", franchise, customer_type, region, level, drill_region, drill_city)
        labels = []
        for r in rows24 + rows25:
            if r["label"] not in labels:
                labels.append(r["label"])
        v24 = {r["label"]: r["sales"] for r in rows24}
        v25 = {r["label"]: r["sales"] for r in rows25}
        b25 = {r["label"]: r["budget"] for r in rows25}
        merged = [{"label": l, "y2024": v24.get(l, 0), "y2025": v25.get(l, 0),
                   "budget2025": b25.get(l, 0)} for l in labels]
        merged.sort(key=lambda d: d["y2024"] + d["y2025"], reverse=True)
        return merged[:15]

    @http.route('/pbi_dashboards/sales/region_drill', type='json', auth='user')
    def region_drill(self, year='2025', franchise='Midea', customerType='all', region='all',
                      level='region', drillRegion=None, drillCity=None):
        if not self._has_access():
            return {'error': 'You do not have access to this dashboard.'}
        try:
            data = self._fetch_region_drill(year, franchise, customerType, region, level, drillRegion, drillCity)
            return {'data': data}
        except Exception as e:
            return {'error': str(e)}

    # ------------------------------------------------------------------
    # access control — mirrors the menu's own hidden-until-granted gate
    # ------------------------------------------------------------------
    def _has_access(self):
        user = request.env.user
        menu = request.env.ref("pbi_dashboards.menu_pbi_sales_dashboard_app", raise_if_not_found=False)
        if not menu:
            return False
        allowed = request.env["dashboard.rights.menu"].sudo().allowed_menu_ids(user)
        return menu.id in allowed

    def _fetch_notes_map(self, year, franchise, customer_type='all', region='all'):
        notes = request.env['pbi.dashboard.note'].sudo().search([
            ('year', '=', year), ('franchise', '=', franchise),
            ('customer_type', '=', customer_type), ('region', '=', region),
        ])
        return {n.key: n.text for n in notes}

    def _current_narrative_texts(self, year, franchise, customer_type='all', region='all'):
        """Recomputes the same narratives _build_pptx would write into a fresh
        export's speaker notes right now, in the same **bold**/##heading##
        marked form _pptx_read_notes_marked extracts from an uploaded file —
        used by import_notes to tell an actual hand edit apart from notes the
        user left untouched."""
        compare = year == 'compare'
        bundle = (self._fetch_compare(franchise, customer_type, region) if compare
                  else self._fetch_bundle(year, franchise, customer_type, region))

        if compare:
            texts = {
                'trend': _pnarrative_trend(bundle['narrative'], franchise),
                'region': _pnarrative_category_compare(bundle['region'], 'regions'),
                'franchise': _pnarrative_category_compare(bundle['franchise'], 'franchises'),
                'customerType': _pnarrative_category_compare(bundle['customerType'], 'customer types'),
                'product': _pnarrative_product_compare(bundle['product']),
                'salesman': _pnarrative_salesman_compare(bundle['salesman']),
            }
        else:
            kpis = bundle['kpis']
            has_prev_year = kpis.get('prevYearSales') is not None
            is_partial_month = bool(kpis.get('isPartialMonth'))

            def dim(key, label, value_key='sales'):
                data = bundle[key]
                if has_prev_year:
                    return _pnarrative_category_single3(data, label, is_partial_month, value_key)
                return _pnarrative_category_single(data, label, value_key)

            product = (_pnarrative_product_single3(bundle['product'], is_partial_month) if has_prev_year
                       else _pnarrative_product_single(bundle['product']))
            salesman = (_pnarrative_salesman_single3(bundle['salesman'], is_partial_month) if has_prev_year
                        else _pnarrative_salesman_single(bundle['salesman']))
            texts = {
                'trend': _pnarrative_trend(bundle['narrative'], franchise),
                'region': dim('region', 'regions'),
                'franchise': dim('franchise', 'franchises'),
                'customerType': dim('customerType', 'customer types'),
                'product': product,
                'salesman': salesman,
            }

        return {key: _pptx_marked_text(text) for key, text in texts.items()}

    @http.route('/pbi_dashboards/sales/data', type='json', auth='user')
    def sales_dashboard_data(self, year='2025', franchise='Midea', customerType='all', region='all'):
        if not self._has_access():
            return {'error': 'You do not have access to this dashboard.'}
        try:
            data = (self._fetch_compare(franchise, customerType, region) if year == 'compare'
                    else self._fetch_bundle(year, franchise, customerType, region))
            data['notes'] = self._fetch_notes_map(year, franchise, customerType, region)
            return data
        except Exception as e:
            return {'error': str(e)}

    @http.route('/pbi_dashboards/sales/import_notes', type='http', auth='user', methods=['POST'])
    def import_notes(self, year='2025', franchise='Midea', customerType='all', region='all', **kwargs):
        def _json(payload, status=200):
            return request.make_response(
                json.dumps(payload), headers=[('Content-Type', 'application/json')], status=status)

        if not self._has_access():
            return _json({'error': 'You do not have access to this dashboard.'}, status=403)

        upload = request.httprequest.files.get('file')
        if not upload:
            return _json({'error': 'No file uploaded.'}, status=400)

        try:
            prs = Presentation(BytesIO(upload.read()))
        except Exception as e:
            return _json({'error': f'Could not read PowerPoint file: {e}'}, status=400)

        # The file's own embedded period/franchise/customerType/region (set at
        # export time, see _build_pptx) is authoritative — the browser's
        # currently-selected filters are just a fallback for files exported
        # before this existed, since the user may well have switched screens
        # since exporting.
        keywords = prs.core_properties.keywords or ''
        m = re.match(
            r'^pbi_dashboards\|year=(?P<year>[^|]*)\|franchise=(?P<franchise>[^|]*)'
            r'\|customerType=(?P<customerType>[^|]*)\|region=(?P<region>.*)$', keywords)
        if m:
            year = m.group('year')
            franchise = m.group('franchise')
            customerType = m.group('customerType')
            region = m.group('region')

        notes_map = {}
        for slide in prs.slides:
            key = _pptx_note_section_key(_pptx_slide_title(slide))
            if not key or not slide.has_notes_slide:
                continue
            text = _pptx_read_notes_marked(slide.notes_slide.notes_text_frame).strip()
            if text:
                notes_map[key] = text

        try:
            baseline = self._current_narrative_texts(year, franchise, customerType, region)
        except Exception:
            baseline = {}

        Note = request.env['pbi.dashboard.note'].sudo()
        applied = {}
        for key, text in notes_map.items():
            existing = Note.search([
                ('key', '=', key), ('year', '=', year), ('franchise', '=', franchise),
                ('customer_type', '=', customerType), ('region', '=', region),
            ], limit=1)
            # Text that still matches what the dashboard would generate right
            # now was never actually hand-edited — don't treat it as an
            # override, and clear a stale one if the user reverted their edit.
            if text.strip() == (baseline.get(key) or '').strip():
                if existing:
                    existing.unlink()
                continue
            applied[key] = text
            if existing:
                existing.text = text
            else:
                Note.create({'key': key, 'year': year, 'franchise': franchise,
                              'customer_type': customerType, 'region': region, 'text': text})

        return _json({'year': year, 'franchise': franchise, 'customerType': customerType, 'region': region,
                       'notes': applied})
