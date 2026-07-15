import json
import logging
import re
from io import BytesIO

from odoo import http
from odoo.http import request

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn, nsdecls
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from .main import (
    _pptx_read_notes_marked, _pptx_marked_text, _pptx_fill_marked_text,
    _pfmt_m, _pfmt, _pshare, MONTH_NAMES,
)

_logger = logging.getLogger(__name__)

MONTH_ABBR = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]

PPTX_COLOR_2024 = RGBColor(0x2a, 0x78, 0xd6)
PPTX_COLOR_2025 = RGBColor(0x1b, 0xaf, 0x7a)
PPTX_COLOR_BUDGET = RGBColor(0xed, 0xa1, 0x00)
PPTX_PIE_PALETTE = [RGBColor(0x2a, 0x78, 0xd6), RGBColor(0x1b, 0xaf, 0x7a), RGBColor(0xed, 0xa1, 0x00),
                    RGBColor(0x4a, 0x3a, 0xa7), RGBColor(0x17, 0xa2, 0xb8), RGBColor(0xe3, 0x49, 0x48),
                    RGBColor(0x8e, 0x44, 0xad), RGBColor(0x16, 0xa0, 0x85), RGBColor(0xd3, 0x54, 0x00),
                    RGBColor(0x2c, 0x3e, 0x50)]
# Excel/PowerPoint custom number-format codes for the value axis — each
# trailing comma divides the raw value by 1000, matching the web
# dashboard's fmtM ("630.6M")/fmtK ("22.8K") abbreviations instead of
# python-pptx's default "General" (raw unformatted numbers).
PPTX_AXIS_NUMBER_FORMAT = {'M': '#,##0.0,,"M"', 'K': '#,##0.0,"K"'}

# Slide geometry for a "page" slide (13.333 x 7.5in) — mirrors the web
# page's `.chart-row` grid (chart-col: 1fr, .narrative: 340px, out of a
# 1240px-max .wrap) by carving the same ~27% off the right for a visible
# notes column instead of hiding that text in PowerPoint's Notes pane only.
PPTX_MARGIN = 0.4
PPTX_CONTENT_TOP = 0.95
PPTX_CONTENT_H = 6.35
PPTX_GRID_GAP = 0.2
PPTX_NOTES_W = 3.4
PPTX_NOTES_GAP = 0.25
PPTX_CHART_AREA_W = 13.333 - 2 * PPTX_MARGIN - PPTX_NOTES_W - PPTX_NOTES_GAP
PPTX_NOTES_LEFT = PPTX_MARGIN + PPTX_CHART_AREA_W + PPTX_NOTES_GAP
PPTX_FULL_W = 13.333 - 2 * PPTX_MARGIN

# KPI tile colors — matches .tile-color-0/1/2 in sales_mail_dashboard.css
# (--series-1 blue, #16a2a6 teal, --series-3 orange), cycled This Year/
# Target/Last Year per renderKpis()'s `tile-color-${i % 3}` in
# sales_mail_dashboard_new.js. The two "good"/"bad" badge colors are that
# same CSS's `[class*="tile-color-"] .sub .good/.bad` (light tints, since
# the tile background is already a saturated color).
PPTX_TILE_COLORS = [RGBColor(0x2a, 0x78, 0xd6), RGBColor(0x16, 0xa2, 0xa6), RGBColor(0xed, 0xa1, 0x00)]
PPTX_TILE_GOOD = RGBColor(0xb8, 0xff, 0xc9)
PPTX_TILE_BAD = RGBColor(0xff, 0xd2, 0xd2)

# comboBarLineChart's secondary-axis line colors (--series-5 violet /
# --series-6 red in sales_dashboard.css's lineColors()) — Page 2 only.
PPTX_LINE_ACHIEVEMENT = RGBColor(0x4a, 0x3a, 0xa7)
PPTX_LINE_YOY = RGBColor(0xe3, 0x49, 0x48)

# Base scope for every Amount/Qty aggregate in this dashboard — all 3
# conditions are fixed/hardcoded per explicit instruction, not tied to the
# Franchise selector: customer number pattern, franchise code, and the
# product-family whitelist (confirmed against live dbprod data). Because
# franchise is hardcoded to MDA here, selecting BEKO/CANDY in the Franchise
# filter will correctly return empty data (bi_franchisecode='MDA' AND
# bi_franchisename='BEKO' can never both be true) — this dashboard's
# formula is Midea-specific by design, matching the source spec.
PRODUCT_SCOPE_CODES = ['ACACC', 'ACAIP', 'ACCON', 'ACCST', 'ACHCL', 'ACPAC', 'ACPKG', 'ACPOR', 'ACVRF', 'ACWIN', 'ACWTS', 'ATOM']
# The literal 'V%' LIKE-pattern must be escaped as 'V%%' — psycopg2 scans for
# %s-style tokens whenever params are passed, so a bare % anywhere in the SQL
# text (not just in actual placeholders) breaks its substitution parser.
BASE_SCOPE_SQL = "bi_cstno NOT LIKE 'V%%' AND bi_franchisecode = 'MDA' AND bi_pgroupcode = ANY(%s)"
BASE_SCOPE_PARAMS = [PRODUCT_SCOPE_CODES]

# Amount = qty * (invprice - invdisc - invcstspldisc), per explicit
# instruction — confirmed necessary against the real bidata export: raw
# bi_amount is populated on only ~14% of in-scope rows, while bi_invprice
# covers ~68%, so bi_amount alone silently under-counts actual sales.
# bi_budgetamount (Target) is a separate, well-populated feed and is
# deliberately NOT run through this formula — only actual This-Year/Last-
# Year sales use it. The qty inside this formula is intentionally
# ungated by _QTY_GATE_SQL (that gate only applies to the standalone Qty
# measure, not this price-component usage). bi_qty itself must be
# COALESCEd too: some sub-groups (e.g. CONG001, ELITE PLUS R22) have
# bi_qty NULL on every in-scope row, which made this expression evaluate
# to NULL rather than 0 — and since sum() over all-NULL inputs returns
# NULL, every "ORDER BY <sales> DESC" (top-N breakdowns) sorted those
# NULL groups AHEAD of real, large sales figures (Postgres's DESC default
# is NULLS FIRST), silently corrupting every Top-N chart's ranking.
AMOUNT_EXPR = "(COALESCE(bi_qty, 0) * (COALESCE(bi_invprice, 0) - COALESCE(bi_invdisc, 0) - COALESCE(bi_invcstspldisc, 0)))"

# Quantity-only gate: catalogflags.cat_flag = '02' marks countable-unit
# parts. A given cat_part can carry several catalogflags rows simultaneously
# (confirmed: e.g. MI-112T2/DHN1-BA5 has both flag '10' and flag '02' rows
# on the same date, not a supersession history) — an EXISTS check is the
# correct semantic ("does this part carry a '02' flag"), not a deduped
# join. Amount is deliberately NOT gated by this (see _qty_gate/
# BASE_SCOPE_SQL split below) — only the actual qty/prevYearQty CASE WHEN
# branches get the extra "AND EXISTS (...)". budgetQty (Target) is
# deliberately NEVER gated by this either, even where qty_gate is threaded
# through: confirmed every bi_budgetqty>0 row across the whole table has an
# empty bi_invpartno (targets are set at a coarser grain than individual
# parts), so gating it the same way would silently zero out every Target
# Qty figure.
_QTY_GATE_SQL = """
    AND EXISTS (SELECT 1 FROM catalogflags cgf WHERE cgf.cat_part = bidata.bi_invpartno AND cgf.cat_flag = '02')
"""

FRANCHISE_OPTIONS = [
    {'v': 'all', 'l': 'All'}, {'v': 'Midea', 'l': 'Midea'}, {'v': 'BEKO', 'l': 'BEKO'}, {'v': 'CANDY', 'l': 'CANDY'},
]

# Region taxonomy (verified against live dbprod data): bi_cstregioncode only
# has 3 values (10 Central/20 Western/30 Eastern), with Riyadh/Qassim living
# inside Central via bi_cstsubregioncode (RYD/QAS). None of that matches
# t_rptregionsdesc.rr_code (100/200/300/400/500) directly, so this maps to
# the matching rr_code first, then looks up the formal desc text. Southern
# (500) has no source data in bidata at all — it simply won't populate.
REGION_EXPR = """
    (SELECT rr_desc FROM t_rptregionsdesc WHERE rr_lang = 1 AND rr_code = (
        CASE WHEN bi_cstsubregioncode = 'RYD' THEN '200'
             WHEN bi_cstsubregioncode = 'QAS' THEN '100'
             WHEN bi_cstregioncode = '20' THEN '300'
             WHEN bi_cstregioncode = '30' THEN '400'
             ELSE NULL END
    ) LIMIT 1)
"""
REGION_FIXED = ['Qassim Region', 'Riyadh Region', 'Western Region', 'Eastern Region']
# Single-region filter clauses, keyed the same as REGION_LABELS — used by
# every per-region-toggle page (5, 15, 16). Still keyed on bidata's own raw
# region/subregion codes, unaffected by the rr_desc lookup above.
REGION_SUBQUERY = {
    'riyadh': ("bi_cstsubregioncode = %s", 'RYD'),
    'qassim': ("bi_cstsubregioncode = %s", 'QAS'),
    'western': ("bi_cstregioncode = %s", '20'),
    'eastern': ("bi_cstregioncode = %s", '30'),
}
REGION_LABELS = {'riyadh': 'Riyadh Region', 'qassim': 'Qassim Region', 'western': 'Western Region', 'eastern': 'Eastern Region'}
# Same 4-way split as REGION_SUBQUERY, but as a single SQL expression
# outputting the slug text directly — lets a region-toggle breakdown be
# grouped by region in one query instead of one query per REGION_SUBQUERY
# entry (see _breakdown_ty_ly_multi/_breakdown_by_quarter_multi).
REGION_SLUG_EXPR = """
    CASE WHEN bi_cstsubregioncode = 'RYD' THEN 'riyadh'
         WHEN bi_cstsubregioncode = 'QAS' THEN 'qassim'
         WHEN bi_cstregioncode = '20' THEN 'western'
         WHEN bi_cstregioncode = '30' THEN 'eastern'
         ELSE NULL END
"""

# Channel/department label source: bi_csttypecode resolved against
# ``t_cstclasstypedesc`` (cs_code/cs_desc, cs_lang='1' for English) — the
# same table/expression sales_kpi_main.py's CUSTOMER_TYPE_EXPR uses for its
# own top drill level (duplicated here, not imported, since sales_kpi_main.py
# already imports FROM this module — importing back would be circular).
# An earlier revision of this file switched Department to a t_salpurgroup-
# based expression instead (mapping bi_csttypecode's 01-04 to Dealers/
# Projects/Key Accounts/Corporate Sales) — that mislabeled every non-Dealers
# department (e.g. code '04', real "Projects" per t_cstclasstypedesc and per
# this dashboard's own Page 17 title "Projects — Product Groups" — see
# projects_productgroups' hardcoded bi_csttypecode = '04' filter below —
# rendered as "Corporate Sales" instead), the same mislabeling
# sales_kpi_main.py's own CUSTOMER_TYPE_EXPR docstring already documented
# for its dashboard. Reverted back to t_cstclasstypedesc so the two
# dashboards agree and Page 17's title matches its own data again. The
# internal toggle slugs (dealers/mt/ws/projects) don't need renaming —
# t_cstclasstypedesc's 02/03/04 are Modern Trade/Whole Sale/Projects, which
# is what those slugs already meant before the t_salpurgroup detour.
CHANNEL_DIM_EXPR = """
    COALESCE((SELECT cs_desc FROM t_cstclasstypedesc WHERE cs_lang = '1' AND cs_code = bi_csttypecode LIMIT 1), 'Others')
"""
CHANNEL_FIXED_CODES = [('01', 'Dealers'), ('02', 'Modern Trade'), ('03', 'Whole Sale'), ('04', 'Projects')]
CHANNEL_CODES = {'dealers': '01', 'mt': '02', 'ws': '03', 'projects': '04'}
CHANNEL_LABELS = {'dealers': 'Dealers', 'mt': 'Modern Trade', 'ws': 'Whole Sale', 'projects': 'Projects'}

# "Main Group" has no real bi_mpdesc-style column on bidata — every
# bi_mpdesc-keyed page uses this bi_pgroupcode-derived CASE mapping to 5
# buckets instead.
MAIN_GROUP_FIXED = ['Windows', 'Split', 'CAC', 'LCAC', 'Others']
MAIN_GROUP_EXPR_DERIVED = """
    CASE WHEN bi_pgroupcode = 'ACWIN' THEN 'Windows'
         WHEN bi_pgroupcode = 'ACWTS' THEN 'Split'
         WHEN bi_pgroupcode = 'ACVRF' THEN 'CAC'
         WHEN bi_pgroupcode IN ('ACCST', 'ACCON', 'ACPAC', 'ACPKG', 'ACPOR') THEN 'LCAC'
         ELSE 'Others' END
"""

# Slide-title-substring -> narrative/notes key, mirroring
# PPTX_NOTE_SECTION_MATCHERS in main.py so the notes importer matches by
# content, not slide position. Several titles contain more than one needle
# (e.g. "Dealers by Region — Main Groups — Riyadh" contains both "main
# groups" and its own "dealers by region — main groups") — resolution picks
# whichever matching needle is LONGEST (see _pptx_note_section_key), so
# these don't need to be hand-ordered by specificity.
PPTX_MAIL_SECTION_MATCHERS = (
    ("total company", "total-company"),
    ("department & region performance", "dept-region-trends"),
    ("quarterly sales progression — company", "quarterly-company"),
    ("quarterly progression by department", "quarterly-department"),
    ("quarterly progression by region", "quarterly-region"),
    ("main groups distribution", "maingroups-distribution"),
    ("main groups", "main-groups"),
    ("sub-groups — top 8", "subgroups-top8"),
    ("sub-groups distribution", "subgroups-distribution"),
    ("departments distribution", "departments-distribution"),
    ("departments", "departments"),
    ("lcac sub-groups", "lcac-subgroups"),
    ("top-3 sub-groups", "top3-by-channel"),
    ("regions by channel", "regions-by-channel"),
    ("dealers by region — main groups", "dealers-region-maingroups"),
    ("dealers by region — split sub-groups", "dealers-region-subgroups"),
    ("projects — product groups", "projects-productgroups"),
)


def _pptx_slide_title(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape.text_frame.text.strip()
    return ""


def _pptx_find_notes_sidebar(slide):
    """The visible notes textbox _pptx_add_notes_sidebar adds to every
    content slide, identified by its fixed left position (PPTX_NOTES_LEFT)
    — the only shape pinned there. Notes live ONLY here now (no speaker
    notes are written), so _import_notes reads this instead of
    slide.notes_slide.notes_text_frame."""
    target = Inches(PPTX_NOTES_LEFT)
    for shape in slide.shapes:
        if shape.has_text_frame and abs(shape.left - target) < Inches(0.05):
            return shape
    return None



# Pages 4/5/13/14/15/16 are one-slide-PER-TAB (see _build_pptx) — their
# title ends with " — <tab label>" (e.g. "...Department — Modern Trade",
# "...Main Groups (YTD) — Riyadh Region"). Maps a matched base key to the
# label set its tabs are drawn from, so _pptx_note_section_key can resolve
# that suffix into the SAME per-state key (e.g. 'quarterly-department_mt')
# the frontend's putNarrative() reads (see sales_mail_dashboard_new.js) —
# the base key alone corresponds to nothing displayed on screen per-panel.
PPTX_CHANNEL_TAB_KEYS = {'quarterly-department', 'top3-by-channel', 'regions-by-channel'}
PPTX_REGION_TAB_KEYS = {'quarterly-region', 'dealers-region-maingroups', 'dealers-region-subgroups'}


def _pptx_note_section_key(title):
    lowered = title.lower()
    matches = [(needle, key) for needle, key in PPTX_MAIL_SECTION_MATCHERS if needle in lowered]
    if not matches:
        return None
    # Prefer the longest (most specific) matching needle — e.g. a slide
    # titled "Dealers by Region — Main Groups" matches both the generic
    # "main groups" needle and its own dedicated, longer one; the longer
    # one wins rather than whichever happens to sit first in the matcher
    # list above.
    _, key = max(matches, key=lambda pair: len(pair[0]))
    labels = (CHANNEL_LABELS if key in PPTX_CHANNEL_TAB_KEYS
              else REGION_LABELS if key in PPTX_REGION_TAB_KEYS else None)
    if labels:
        for slug, label in labels.items():
            if title.endswith(label):
                return f'{key}_{slug}'
    return key


class PbiSalesMailController(http.Controller):
    """Backs "PBI Dashboards > Sales Dashboards > Sales Mail Content" — a
    17-page report-style dashboard rebuilt from the client's actual Power BI
    source file (HHS_SalesData_Analysis.pbix), reading live from ``bidata``
    directly (joined to ``t_rptregionsdesc``/``t_cstclasstypedesc`` for region
    and channel/department display text, and ``catalogflags`` for the
    qty-only cat_flag gate — see
    BASE_SCOPE_SQL/REGION_EXPR/CHANNEL_DIM_EXPR/_QTY_GATE_SQL).
    Period (Year+Month) and Franchise filters only; every other dimension
    (channel, region, department toggle) is an in-page toggle, not a global
    filter. No drill-down.
    """

    # ------------------------------------------------------------------
    # low-level helpers
    # ------------------------------------------------------------------
    def _query(self, sql, params=()):
        request.env.cr.execute(sql, params)
        cols = [d[0] for d in request.env.cr.description]
        return [dict(zip(cols, row)) for row in request.env.cr.fetchall()]

    def _has_access_for(self, menu_xmlid):
        user = request.env.user
        menu = request.env.ref(menu_xmlid, raise_if_not_found=False)
        if not menu:
            return False
        allowed = request.env["dashboard.rights.menu"].sudo().allowed_menu_ids(user)
        return menu.id in allowed

    def _has_access(self):
        return self._has_access_for("pbi_dashboards.menu_pbi_sales_mail_dashboard")

    # Backs "Sales Analysis - New" — same bidata-direct data/PPTX export/
    # notes-import as "Sales Analysis" above (only the on-screen Prev/Next
    # pagination differs, see sales_mail_dashboard_new.js), gated by its own
    # menu grant (menu_pbi_sales_mail_dashboard_new) rather than reusing
    # "Sales Analysis"'s, per this module's one-grant-per-leaf convention.
    def _has_access_new(self):
        return self._has_access_for("pbi_dashboards.menu_pbi_sales_mail_dashboard_new")

    def _safe(self, label, fn, default):
        """One page's query failing shouldn't 500 the entire dashboard —
        every page block in _fetch_bundle is wrapped in this so the rest of
        the report still renders."""
        try:
            return fn()
        except Exception as e:
            _logger.warning("sales_mail_main: %s failed: %s", label, e)
            # A failed query (e.g. bi_mpdesc missing) leaves the shared
            # request cursor's transaction aborted — every later query would
            # fail too with "current transaction is aborted" unless rolled
            # back here, even ones that don't touch the offending column.
            request.env.cr.rollback()
            return default

    def _qty_gate(self):
        """Probed once per request: returns _QTY_GATE_SQL (the catalogflags
        cat_flag='02' EXISTS check) when catalogflags actually has '02' rows,
        else '' (no gate) — a defensive fallback in case catalogflags is ever
        empty again, matching the same probe pattern used elsewhere in this
        controller rather than assuming the table is always populated."""
        try:
            rows = self._query("SELECT 1 FROM catalogflags WHERE cat_flag = '02' LIMIT 1")
            return _QTY_GATE_SQL if rows else ""
        except Exception:
            request.env.cr.rollback()
            return ""

    # ------------------------------------------------------------------
    # filter options
    # ------------------------------------------------------------------
    def _year_options(self):
        rows = self._query(f"""
            SELECT DISTINCT bi_year FROM bidata
            WHERE bi_qty IS NOT NULL AND bi_year IS NOT NULL AND {BASE_SCOPE_SQL}
            ORDER BY 1 DESC LIMIT 8
        """, BASE_SCOPE_PARAMS)
        return [{'v': r['bi_year'], 'l': str(r['bi_year'])} for r in rows]

    def _latest_period(self):
        rows = self._query(f"""
            SELECT bi_year, bi_month FROM bidata
            WHERE bi_qty IS NOT NULL AND bi_year IS NOT NULL AND bi_month IS NOT NULL
                  AND {BASE_SCOPE_SQL}
            ORDER BY bi_year DESC, bi_month DESC LIMIT 1
        """, BASE_SCOPE_PARAMS)
        return (rows[0]['bi_year'], rows[0]['bi_month']) if rows else (None, None)

    # ------------------------------------------------------------------
    # filter/period clauses
    # ------------------------------------------------------------------
    def _dim_clauses(self, franchise):
        if franchise and franchise != 'all':
            return ["bi_franchisename = %s"], [franchise]
        return [], []

    @staticmethod
    def _period_window(month, mode):
        """mode: 'month' | 'qtd' | 'ytd' -> (month_gte, month_lte), inclusive."""
        if mode == 'month':
            return month, month
        if mode == 'qtd':
            qstart = ((month - 1) // 3) * 3 + 1
            return qstart, month
        return 1, month  # ytd

    # ------------------------------------------------------------------
    # aggregates — conditional aggregation pulls TY+LY in one round trip.
    # ------------------------------------------------------------------
    def _sums_ty_ly(self, year, month, mode, franchise, extra_clauses=(), extra_params=(), qty_gate=''):
        gte, lte = self._period_window(month, mode)
        dim_clauses, dim_params = self._dim_clauses(franchise)
        years = [year, year - 1] if (year - 1) >= 2023 else [year]
        clauses = ["bi_month BETWEEN %s AND %s", "bi_year = ANY(%s)", BASE_SCOPE_SQL] + dim_clauses + list(extra_clauses)
        params = [gte, lte, years] + BASE_SCOPE_PARAMS + dim_params + list(extra_params)
        where = " AND ".join(clauses)
        r = self._query(f"""
            SELECT
              sum(CASE WHEN bi_year = %s THEN {AMOUNT_EXPR} ELSE 0 END) AS sales,
              sum(CASE WHEN bi_year = %s THEN bi_budgetamount ELSE 0 END) AS budget,
              sum(CASE WHEN bi_year = %s {qty_gate} THEN COALESCE(bi_qty, 0) ELSE 0 END) AS qty,
              sum(CASE WHEN bi_year = %s THEN bi_budgetqty ELSE 0 END) AS budget_qty,
              sum(CASE WHEN bi_year = %s THEN {AMOUNT_EXPR} ELSE 0 END) AS prev_sales,
              sum(CASE WHEN bi_year = %s {qty_gate} THEN COALESCE(bi_qty, 0) ELSE 0 END) AS prev_qty
            FROM bidata WHERE {where}
        """, [year, year, year, year, year - 1, year - 1] + params)[0]
        return {
            'sales': float(r['sales'] or 0), 'budget': float(r['budget'] or 0),
            'qty': int(r['qty'] or 0), 'budgetQty': int(r['budget_qty'] or 0),
            'prevYearSales': float(r['prev_sales'] or 0), 'prevYearQty': int(r['prev_qty'] or 0),
        }

    def _breakdown_ty_ly(self, year, month, mode, franchise, dim_expr, extra_clauses=(), extra_params=(),
                          fixed_labels=None, limit=20, qty_gate=''):
        dim_clauses, dim_params = self._dim_clauses(franchise)
        gte, lte = self._period_window(month, mode)
        years = [year, year - 1] if (year - 1) >= 2023 else [year]
        clauses = ["bi_month BETWEEN %s AND %s", "bi_year = ANY(%s)", BASE_SCOPE_SQL, f"({dim_expr}) IS NOT NULL"]
        clauses += dim_clauses + list(extra_clauses)
        params = [gte, lte, years] + BASE_SCOPE_PARAMS + dim_params + list(extra_params)
        where = " AND ".join(clauses)
        rows = self._query(f"""
            SELECT {dim_expr} AS label,
              sum(CASE WHEN bi_year = %s THEN {AMOUNT_EXPR} ELSE 0 END) AS sales,
              sum(CASE WHEN bi_year = %s THEN bi_budgetamount ELSE 0 END) AS budget,
              sum(CASE WHEN bi_year = %s {qty_gate} THEN COALESCE(bi_qty, 0) ELSE 0 END) AS qty,
              sum(CASE WHEN bi_year = %s THEN bi_budgetqty ELSE 0 END) AS budget_qty,
              sum(CASE WHEN bi_year = %s THEN {AMOUNT_EXPR} ELSE 0 END) AS prev_sales,
              sum(CASE WHEN bi_year = %s {qty_gate} THEN COALESCE(bi_qty, 0) ELSE 0 END) AS prev_qty
            FROM bidata WHERE {where}
            GROUP BY 1
            ORDER BY 2 DESC
            LIMIT %s
        """, [year, year, year, year, year - 1, year - 1] + params + [limit])
        by_label = {r['label']: r for r in rows}

        def _row(label, r):
            return {
                'code': label, 'label': label,
                'sales': float(r['sales'] or 0) if r else 0.0,
                'budget': float(r['budget'] or 0) if r else 0.0,
                'qty': int(r['qty'] or 0) if r else 0,
                'budgetQty': int(r['budget_qty'] or 0) if r else 0,
                'prevYearSales': float(r['prev_sales'] or 0) if r else 0.0,
                'prevYearQty': int(r['prev_qty'] or 0) if r else 0,
            }

        if fixed_labels:
            # zero-fill every fixed category (so one with no activity this
            # period still gets a 0-value bar instead of vanishing), then
            # sort by This-Year sales descending like every other breakdown.
            zero_filled = [_row(label, by_label.get(label)) for label in fixed_labels]
            zero_filled.sort(key=lambda r: r['sales'], reverse=True)
            return zero_filled
        return [_row(r['label'], r) for r in rows]

    def _breakdown_ty_ly_multi(self, year, month, mode, franchise, group_expr, group_keys, dim_expr,
                                fixed_labels, extra_clauses=(), extra_params=(), qty_gate=''):
        """Same shape/output as _breakdown_ty_ly(fixed_labels=...) called
        once per `group_keys` entry, but as ONE query grouped by both
        `group_expr` (e.g. channel code or region slug) and `dim_expr`
        instead of one query per group — used by the 4-toggle-state pages
        (regions-by-channel, dealers-region-maingroups) to cut 4 round
        trips down to 1. Returns {group_key: [dim_expr rows]}, each list
        zero-filled/sorted exactly like the single-group version."""
        dim_clauses, dim_params = self._dim_clauses(franchise)
        gte, lte = self._period_window(month, mode)
        years = [year, year - 1] if (year - 1) >= 2023 else [year]
        clauses = ["bi_month BETWEEN %s AND %s", "bi_year = ANY(%s)", BASE_SCOPE_SQL,
                   f"({group_expr}) IS NOT NULL", f"({dim_expr}) IS NOT NULL"]
        clauses += dim_clauses + list(extra_clauses)
        params = [gte, lte, years] + BASE_SCOPE_PARAMS + dim_params + list(extra_params)
        where = " AND ".join(clauses)
        rows = self._query(f"""
            SELECT {group_expr} AS grp, {dim_expr} AS label,
              sum(CASE WHEN bi_year = %s THEN {AMOUNT_EXPR} ELSE 0 END) AS sales,
              sum(CASE WHEN bi_year = %s THEN bi_budgetamount ELSE 0 END) AS budget,
              sum(CASE WHEN bi_year = %s {qty_gate} THEN COALESCE(bi_qty, 0) ELSE 0 END) AS qty,
              sum(CASE WHEN bi_year = %s THEN bi_budgetqty ELSE 0 END) AS budget_qty,
              sum(CASE WHEN bi_year = %s THEN {AMOUNT_EXPR} ELSE 0 END) AS prev_sales,
              sum(CASE WHEN bi_year = %s {qty_gate} THEN COALESCE(bi_qty, 0) ELSE 0 END) AS prev_qty
            FROM bidata WHERE {where}
            GROUP BY 1, 2
        """, [year, year, year, year, year - 1, year - 1] + params)

        by_grp = {}
        for r in rows:
            by_grp.setdefault(r['grp'], {})[r['label']] = r

        def _row(label, r):
            return {
                'code': label, 'label': label,
                'sales': float(r['sales'] or 0) if r else 0.0,
                'budget': float(r['budget'] or 0) if r else 0.0,
                'qty': int(r['qty'] or 0) if r else 0,
                'budgetQty': int(r['budget_qty'] or 0) if r else 0,
                'prevYearSales': float(r['prev_sales'] or 0) if r else 0.0,
                'prevYearQty': int(r['prev_qty'] or 0) if r else 0,
            }

        out = {}
        for grp in group_keys:
            by_label = by_grp.get(grp, {})
            zero_filled = [_row(label, by_label.get(label)) for label in fixed_labels]
            zero_filled.sort(key=lambda r: r['sales'], reverse=True)
            out[grp] = zero_filled
        return out

    def _breakdown_ty_ly_multi_topn(self, year, month, mode, franchise, group_expr, group_keys, dim_expr,
                                     limit, extra_clauses=(), extra_params=(), qty_gate=''):
        """Same shape as _breakdown_ty_ly(limit=N) called once per
        `group_keys` entry, but as ONE query — top-N `dim_expr` rows per
        `group_expr` group via ROW_NUMBER(), instead of one full-table
        aggregate query per group. Used by the 4-toggle-state top-N pages
        (top3-by-channel, dealers-region-subgroups), mirroring how
        _breakdown_ty_ly_multi already batches the fixed-label 4-toggle
        pages (regions-by-channel, dealers-region-maingroups). Returns
        {group_key: [dim_expr rows]}, sorted by This-Year sales descending
        like every other breakdown — groups with no matching rows get []
        rather than a KeyError."""
        dim_clauses, dim_params = self._dim_clauses(franchise)
        gte, lte = self._period_window(month, mode)
        years = [year, year - 1] if (year - 1) >= 2023 else [year]
        clauses = ["bi_month BETWEEN %s AND %s", "bi_year = ANY(%s)", BASE_SCOPE_SQL,
                   f"({group_expr}) IS NOT NULL", f"({dim_expr}) IS NOT NULL"]
        clauses += dim_clauses + list(extra_clauses)
        params = [gte, lte, years] + BASE_SCOPE_PARAMS + dim_params + list(extra_params)
        where = " AND ".join(clauses)
        sales_case = f"sum(CASE WHEN bi_year = %s THEN {AMOUNT_EXPR} ELSE 0 END)"
        rows = self._query(f"""
            SELECT grp, label, sales, budget, qty, budget_qty, prev_sales, prev_qty FROM (
                SELECT {group_expr} AS grp, {dim_expr} AS label,
                  {sales_case} AS sales,
                  sum(CASE WHEN bi_year = %s THEN bi_budgetamount ELSE 0 END) AS budget,
                  sum(CASE WHEN bi_year = %s {qty_gate} THEN COALESCE(bi_qty, 0) ELSE 0 END) AS qty,
                  sum(CASE WHEN bi_year = %s THEN bi_budgetqty ELSE 0 END) AS budget_qty,
                  sum(CASE WHEN bi_year = %s THEN {AMOUNT_EXPR} ELSE 0 END) AS prev_sales,
                  sum(CASE WHEN bi_year = %s {qty_gate} THEN COALESCE(bi_qty, 0) ELSE 0 END) AS prev_qty,
                  ROW_NUMBER() OVER (PARTITION BY {group_expr} ORDER BY {sales_case} DESC) AS rn
                FROM bidata WHERE {where}
                GROUP BY 1, 2
            ) t
            WHERE rn <= %s
            ORDER BY grp, sales DESC
        """, [year, year, year, year, year - 1, year - 1, year] + params + [limit])

        out = {grp: [] for grp in group_keys}
        for r in rows:
            out.setdefault(r['grp'], []).append({
                'code': r['label'], 'label': r['label'],
                'sales': float(r['sales'] or 0), 'budget': float(r['budget'] or 0),
                'qty': int(r['qty'] or 0), 'budgetQty': int(r['budget_qty'] or 0),
                'prevYearSales': float(r['prev_sales'] or 0), 'prevYearQty': int(r['prev_qty'] or 0),
            })
        return out

    def _breakdown_by_quarter(self, year, month, franchise, extra_clauses=(), extra_params=(), qty_gate=''):
        """Calendar-year quarterly progression up through the quarter that
        contains the selected month (Q1..current quarter only) — unlike
        every other breakdown here which is scoped to the selected month's
        MTD/QTD/YTD window — used by the 3 quarterly-progression pages.
        This-Year actual sales/qty are capped at the selected month (so
        picking an earlier month doesn't leak later actuals into the
        current quarter); Budget and Last-Year stay full-quarter for
        comparison. Quarters entirely beyond the selected month are omitted
        rather than shown as a misleading 0.0 This-Year bar."""
        dim_clauses, dim_params = self._dim_clauses(franchise)
        years = [year, year - 1] if (year - 1) >= 2023 else [year]
        clauses = ["bi_year = ANY(%s)", BASE_SCOPE_SQL] + dim_clauses + list(extra_clauses)
        params = [years] + BASE_SCOPE_PARAMS + dim_params + list(extra_params)
        where = " AND ".join(clauses)
        rows = self._query(f"""
            SELECT ((bi_month - 1) / 3 + 1) AS quarter,
              sum(CASE WHEN bi_year = %s AND bi_month <= %s THEN {AMOUNT_EXPR} ELSE 0 END) AS sales,
              sum(CASE WHEN bi_year = %s THEN bi_budgetamount ELSE 0 END) AS budget,
              sum(CASE WHEN bi_year = %s AND bi_month <= %s {qty_gate} THEN COALESCE(bi_qty, 0) ELSE 0 END) AS qty,
              sum(CASE WHEN bi_year = %s THEN bi_budgetqty ELSE 0 END) AS budget_qty,
              sum(CASE WHEN bi_year = %s THEN {AMOUNT_EXPR} ELSE 0 END) AS prev_sales,
              sum(CASE WHEN bi_year = %s {qty_gate} THEN COALESCE(bi_qty, 0) ELSE 0 END) AS prev_qty
            FROM bidata WHERE {where}
            GROUP BY 1
        """, [year, month, year, year, month, year, year - 1, year - 1] + params)
        by_q = {r['quarter']: r for r in rows}
        current_quarter = (month - 1) // 3 + 1
        out = []
        for q in range(1, current_quarter + 1):
            r = by_q.get(q)
            out.append({
                'code': f'Q{q}', 'label': f'Q{q}',
                'sales': float(r['sales'] or 0) if r else 0.0,
                'budget': float(r['budget'] or 0) if r else 0.0,
                'qty': int(r['qty'] or 0) if r else 0,
                'budgetQty': int(r['budget_qty'] or 0) if r else 0,
                'prevYearSales': float(r['prev_sales'] or 0) if r else 0.0,
                'prevYearQty': int(r['prev_qty'] or 0) if r else 0,
            })
        return out

    def _breakdown_by_quarter_multi(self, year, month, franchise, group_expr, group_keys,
                                     extra_clauses=(), extra_params=(), qty_gate=''):
        """Same shape/output as _breakdown_by_quarter (including the
        this-year-capped-at-selected-month behavior and clipping to quarters
        up through the one containing the selected month) called once per
        `group_keys` entry, but as ONE query grouped by both `group_expr`
        and quarter — used by quarterly-department/quarterly-region to cut
        4 round trips down to 1. Returns {group_key: [Q1..current-Q rows]}."""
        dim_clauses, dim_params = self._dim_clauses(franchise)
        years = [year, year - 1] if (year - 1) >= 2023 else [year]
        clauses = ["bi_year = ANY(%s)", BASE_SCOPE_SQL, f"({group_expr}) IS NOT NULL"]
        clauses += dim_clauses + list(extra_clauses)
        params = [years] + BASE_SCOPE_PARAMS + dim_params + list(extra_params)
        where = " AND ".join(clauses)
        rows = self._query(f"""
            SELECT {group_expr} AS grp, ((bi_month - 1) / 3 + 1) AS quarter,
              sum(CASE WHEN bi_year = %s AND bi_month <= %s THEN {AMOUNT_EXPR} ELSE 0 END) AS sales,
              sum(CASE WHEN bi_year = %s THEN bi_budgetamount ELSE 0 END) AS budget,
              sum(CASE WHEN bi_year = %s AND bi_month <= %s {qty_gate} THEN COALESCE(bi_qty, 0) ELSE 0 END) AS qty,
              sum(CASE WHEN bi_year = %s THEN bi_budgetqty ELSE 0 END) AS budget_qty,
              sum(CASE WHEN bi_year = %s THEN {AMOUNT_EXPR} ELSE 0 END) AS prev_sales,
              sum(CASE WHEN bi_year = %s {qty_gate} THEN COALESCE(bi_qty, 0) ELSE 0 END) AS prev_qty
            FROM bidata WHERE {where}
            GROUP BY 1, 2
        """, [year, month, year, year, month, year, year - 1, year - 1] + params)

        by_grp_q = {}
        for r in rows:
            by_grp_q.setdefault(r['grp'], {})[r['quarter']] = r

        current_quarter = (month - 1) // 3 + 1
        out = {}
        for grp in group_keys:
            by_q = by_grp_q.get(grp, {})
            qs = []
            for q in range(1, current_quarter + 1):
                r = by_q.get(q)
                qs.append({
                    'code': f'Q{q}', 'label': f'Q{q}',
                    'sales': float(r['sales'] or 0) if r else 0.0,
                    'budget': float(r['budget'] or 0) if r else 0.0,
                    'qty': int(r['qty'] or 0) if r else 0,
                    'budgetQty': int(r['budget_qty'] or 0) if r else 0,
                    'prevYearSales': float(r['prev_sales'] or 0) if r else 0.0,
                    'prevYearQty': int(r['prev_qty'] or 0) if r else 0,
                })
            out[grp] = qs
        return out

    def _append_total_row(self, rows, label='Total'):
        if not rows:
            return rows
        total = {'code': label, 'label': label}
        for key in ('sales', 'budget', 'qty', 'budgetQty', 'prevYearSales', 'prevYearQty'):
            total[key] = sum((r.get(key) or 0) for r in rows)
        return rows + [total]

    def _with_pct(self, rows):
        """Adds achievementPct (TY as % of Target) and yoyPct (TY as % of
        Last Year) to each row — feeds the combo bar+line chart's secondary
        axis, mirroring the pbix's MS_*_Achievement_PC/MS_*_YoY_PC measures."""
        for r in rows:
            r['achievementPct'] = (r['sales'] / r['budget'] * 100) if r.get('budget') else None
            r['yoyPct'] = (r['sales'] / r['prevYearSales'] * 100) if r.get('prevYearSales') else None
        return rows

    # ------------------------------------------------------------------
    # narratives — leader + share + achievement% + YoY%, generic enough to
    # reuse across every breakdown-shaped page here.
    # ------------------------------------------------------------------
    def _narr_breakdown(self, rows, dim_label, value_key='sales', budget_key='budget', prev_key='prevYearSales'):
        rows = [r for r in rows if r.get(value_key) is not None and r.get('label') != 'Total']
        if not rows:
            return f"No data available for {dim_label}."
        top = max(rows, key=lambda d: d.get(value_key) or 0)
        total = sum((d.get(value_key) or 0) for d in rows)
        fmt_fn = _pfmt_m if value_key in ('sales', 'prevYearSales') else _pfmt
        text = (f"{top['label']} leads {dim_label} at {fmt_fn(top.get(value_key, 0))}, "
                f"accounting for {_pshare(top.get(value_key, 0), total)} of the total shown")
        budget = top.get(budget_key)
        if budget:
            achv = top.get(value_key, 0) / budget * 100
            text += f", {'ahead of' if achv >= 100 else 'behind'} target ({achv:.0f}%)"
        prev = top.get(prev_key)
        if prev:
            yoy = (top.get(value_key, 0) - prev) / prev * 100
            text += f", {'up' if yoy >= 0 else 'down'} {abs(yoy):.0f}% year-over-year"
        if len(rows) > 1:
            low = min(rows, key=lambda d: d.get(value_key) or 0)
            if low is not top:
                text += f", while {low['label']} trails at {fmt_fn(low.get(value_key, 0))}"
        return text + "."

    def _narr_total_company(self, mtd_row, ytd_row):
        if not mtd_row or not ytd_row:
            return "No data available."
        parts = []
        for row, label in ((mtd_row, 'This month'), (ytd_row, 'Year-to-date')):
            achv = (row['sales'] / row['budget'] * 100) if row.get('budget') else None
            yoy = ((row['sales'] - row['prevYearSales']) / row['prevYearSales'] * 100) if row.get('prevYearSales') else None
            s = f"{label}: {_pfmt_m(row['sales'])} sales"
            if achv is not None:
                s += f" ({achv:.0f}% of target)"
            if yoy is not None:
                s += f", {'up' if yoy >= 0 else 'down'} {abs(yoy):.0f}% YoY"
            qs = f"{_pfmt(row['qty'])} units"
            qachv = (row['qty'] / row['budgetQty'] * 100) if row.get('budgetQty') else None
            if qachv is not None:
                qs += f" ({qachv:.0f}% qty target)"
            parts.append(f"{s}. {qs}.")
        return "\n".join(parts)

    def _narr_quarterly(self, rows, dim_label):
        rows = [r for r in rows if r.get('sales')]
        if not rows:
            return f"No data available for {dim_label}."
        best = max(rows, key=lambda d: d.get('sales') or 0)
        return f"{best['label']} is the strongest quarter for {dim_label} at {_pfmt_m(best['sales'])}."

    def _combined_narrative(self, narratives, base_key, labels_dict):
        """Concatenates a toggle page's 4 per-state narratives into one
        block (e.g. 'DEALERS\\n...\\n\\nMODERN TRADE\\n...') — these 6 pages
        are single combined slides in the PPTX export (see _build_pptx), so
        there's one narrative per page, not one per state. Shared between
        _fetch_bundle (so this combined key has a baseline import_notes can
        diff a re-imported slide's notes against) and _build_pptx (so the
        exported slide's notes are byte-identical to that baseline before
        any hand edit)."""
        parts = [f"{label.upper()}\n{narratives[f'{base_key}_{slug}']}"
                 for slug, label in labels_dict.items() if narratives.get(f'{base_key}_{slug}')]
        return "\n\n".join(parts)

    # ------------------------------------------------------------------
    # bundle
    # ------------------------------------------------------------------
    def _fetch_bundle(self, year, month, franchise):
        # SET LOCAL scopes to this transaction only (auto-reverts at
        # commit, no persistent/database-wide config change, no effect on
        # other requests sharing the connection pool). Confirmed via
        # EXPLAIN ANALYZE that JIT compilation was costing ~830ms per
        # query on top of ~120ms of actual work — Postgres's cost model
        # treats each of this bundle's ~40 multi-CASE-WHEN aggregate
        # queries as complex enough to cross the JIT threshold even though
        # they scan a modest row count, so the compile overhead never pays
        # for itself. This was the dominant cost behind the whole
        # dashboard taking 14-18s to load.
        request.env.cr.execute("SET LOCAL jit = off")
        has_prev_year = (year - 1) >= 2023
        dim_clauses, dim_params = self._dim_clauses(franchise)
        quarter = (month - 1) // 3 + 1
        month_label = f"{MONTH_NAMES[month - 1]} {year}"
        channel_labels = [l for _, l in CHANNEL_FIXED_CODES]
        mg_expr = MAIN_GROUP_EXPR_DERIVED
        qty_gate = self._qty_gate()

        # KPI tiles — separate MTD/YTD-per-year queries, deliberately
        # matching sales_kpi_main.py's exact _sums pattern.
        def _sums(y, month_eq=None, month_lte=None):
            clauses = ["bi_year = %s", BASE_SCOPE_SQL] + dim_clauses
            params = [y] + BASE_SCOPE_PARAMS + dim_params
            if month_eq is not None:
                clauses.append("bi_month = %s")
                params.append(month_eq)
            if month_lte is not None:
                clauses.append("bi_month <= %s")
                params.append(month_lte)
            where = " AND ".join(clauses)
            r = self._query(f"""
                SELECT sum({AMOUNT_EXPR}) AS sales, sum(bi_budgetamount) AS budget,
                       sum(CASE WHEN true {qty_gate} THEN COALESCE(bi_qty, 0) ELSE 0 END) AS qty,
                       sum(bi_budgetqty) AS budget_qty
                FROM bidata WHERE {where}
            """, params)[0]
            return {'sales': float(r['sales'] or 0), 'budget': float(r['budget'] or 0),
                    'qty': int(r['qty'] or 0), 'budgetQty': int(r['budget_qty'] or 0)}

        mtd_this = _sums(year, month_eq=month)
        ytd_this = _sums(year, month_lte=month)
        if has_prev_year:
            mtd_last = _sums(year - 1, month_eq=month)
            ytd_last = _sums(year - 1, month_lte=month)
        else:
            mtd_last = {'sales': 0.0, 'budget': 0.0, 'qty': 0, 'budgetQty': 0}
            ytd_last = {'sales': 0.0, 'budget': 0.0, 'qty': 0, 'budgetQty': 0}

        kpis = {
            'mtdThisYear': mtd_this['sales'], 'mtdTarget': mtd_this['budget'], 'mtdLastYear': mtd_last['sales'],
            'mtdQtyThisYear': mtd_this['qty'], 'mtdQtyTarget': mtd_this['budgetQty'], 'mtdQtyLastYear': mtd_last['qty'],
            'ytdThisYear': ytd_this['sales'], 'ytdTarget': ytd_this['budget'], 'ytdLastYear': ytd_last['sales'],
            'ytdQtyThisYear': ytd_this['qty'], 'ytdQtyTarget': ytd_this['budgetQty'], 'ytdQtyLastYear': ytd_last['qty'],
        }

        # Page 1: total-company — 4 separate "3-bar, no category" charts.
        def _single_row():
            return [{'code': 'x', 'label': ''}]

        def _total_company_period(mode):
            row = {**self._sums_ty_ly(year, month, mode, franchise, qty_gate=qty_gate), 'code': 'x', 'label': ''}
            return [row]

        total_company = {
            'mtd': self._safe('total_company.mtd', lambda: _total_company_period('month'), _single_row()),
            'ytd': self._safe('total_company.ytd', lambda: _total_company_period('ytd'), _single_row()),
        }

        # Page 2: dept-region-trends — 2 combo charts, + Total row + pct fields.
        dept_trend = self._safe('dept_trend', lambda: self._with_pct(self._append_total_row(
            self._breakdown_ty_ly(year, month, 'ytd', franchise, CHANNEL_DIM_EXPR, fixed_labels=channel_labels,
                                   qty_gate=qty_gate))), [])
        region_trend = self._safe('region_trend', lambda: self._with_pct(self._append_total_row(
            self._breakdown_ty_ly(year, month, 'ytd', franchise, REGION_EXPR,
                                   extra_clauses=["bi_csttypecode = '01'"], fixed_labels=REGION_FIXED,
                                   qty_gate=qty_gate))), [])

        # Page 3: quarterly-company.
        quarterly_company = self._safe('quarterly_company', lambda: self._breakdown_by_quarter(
            year, month, franchise, qty_gate=qty_gate), [])

        # Page 4: quarterly-department (toggle) — one grouped query for all
        # 4 states instead of one query per state (see
        # _breakdown_by_quarter_multi).
        _qdept_by_code = self._safe('quarterly_department', lambda: self._breakdown_by_quarter_multi(
            year, month, franchise, "bi_csttypecode", list(CHANNEL_CODES.values()), qty_gate=qty_gate),
            {code: [] for code in CHANNEL_CODES.values()})
        quarterly_department = {slug: _qdept_by_code.get(code, []) for slug, code in CHANNEL_CODES.items()}

        # Page 5: quarterly-region (toggle, Dealers channel) — likewise one
        # grouped query for all 4 regions.
        quarterly_region = self._safe('quarterly_region', lambda: self._breakdown_by_quarter_multi(
            year, month, franchise, REGION_SLUG_EXPR, list(REGION_SUBQUERY.keys()),
            extra_clauses=["bi_csttypecode = '01'"], qty_gate=qty_gate),
            {slug: [] for slug in REGION_SUBQUERY})

        # Page 6: main-groups — MTD, YTD (derived from bi_pgroupcode).
        main_groups = {
            'mtd': self._safe('main_groups.mtd', lambda: self._breakdown_ty_ly(
                year, month, 'month', franchise, mg_expr, fixed_labels=MAIN_GROUP_FIXED, qty_gate=qty_gate), []),
            'ytd': self._safe('main_groups.ytd', lambda: self._breakdown_ty_ly(
                year, month, 'ytd', franchise, mg_expr, fixed_labels=MAIN_GROUP_FIXED, qty_gate=qty_gate), []),
        }

        # Page 7: subgroups-top8 (YTD) — also feeds page 8's donuts.
        subgroups_top8 = self._safe('subgroups_top8', lambda: self._breakdown_ty_ly(
            year, month, 'ytd', franchise, "bi_psgroupname", limit=8, qty_gate=qty_gate), [])

        # Page 10: departments — MTD, YTD — also feeds page 11's donuts.
        departments = {
            'mtd': self._safe('departments.mtd', lambda: self._breakdown_ty_ly(
                year, month, 'month', franchise, CHANNEL_DIM_EXPR, fixed_labels=channel_labels, qty_gate=qty_gate), []),
            'ytd': self._safe('departments.ytd', lambda: self._breakdown_ty_ly(
                year, month, 'ytd', franchise, CHANNEL_DIM_EXPR, fixed_labels=channel_labels, qty_gate=qty_gate), []),
        }

        # Page 12: lcac-subgroups — MTD, YTD, filtered Main Group = LCAC.
        lcac_subgroups = {
            'mtd': self._safe('lcac_subgroups.mtd', lambda: self._breakdown_ty_ly(
                year, month, 'month', franchise, "bi_psgroupname", extra_clauses=[f"({mg_expr}) = 'LCAC'"],
                qty_gate=qty_gate), []),
            'ytd': self._safe('lcac_subgroups.ytd', lambda: self._breakdown_ty_ly(
                year, month, 'ytd', franchise, "bi_psgroupname", extra_clauses=[f"({mg_expr}) = 'LCAC'"],
                qty_gate=qty_gate), []),
        }

        # Page 13: top3-by-channel (toggle) — one grouped top-N-per-channel
        # query instead of one full-table query per channel state (same
        # batching already applied to pages 14/15 via _breakdown_ty_ly_multi).
        _top3_by_code = self._safe('top3_by_channel', lambda: self._breakdown_ty_ly_multi_topn(
            year, month, 'ytd', franchise, "bi_csttypecode", list(CHANNEL_CODES.values()),
            "bi_psgroupname", limit=3, qty_gate=qty_gate),
            {code: [] for code in CHANNEL_CODES.values()})
        top3_by_channel = {slug: _top3_by_code.get(code, []) for slug, code in CHANNEL_CODES.items()}

        # Page 14: regions-by-channel (toggle) — one grouped query
        # (channel, region) instead of one query per channel state.
        _rbc_by_code = self._safe('regions_by_channel', lambda: self._breakdown_ty_ly_multi(
            year, month, 'ytd', franchise, "bi_csttypecode", list(CHANNEL_CODES.values()),
            REGION_EXPR, REGION_FIXED, qty_gate=qty_gate),
            {code: [] for code in CHANNEL_CODES.values()})
        regions_by_channel = {slug: _rbc_by_code.get(code, []) for slug, code in CHANNEL_CODES.items()}

        # Page 15: dealers-region-maingroups (toggle) — one grouped query
        # (region, main group) instead of one query per region state.
        dealers_region_maingroups = self._safe('dealers_region_maingroups', lambda: self._breakdown_ty_ly_multi(
            year, month, 'ytd', franchise, REGION_SLUG_EXPR, list(REGION_SUBQUERY.keys()),
            mg_expr, MAIN_GROUP_FIXED, extra_clauses=["bi_csttypecode = '01'"], qty_gate=qty_gate),
            {slug: [] for slug in REGION_SUBQUERY})

        # Page 16: dealers-region-subgroups (toggle, Main Group = Split) —
        # one grouped top-8-per-region query (via REGION_SLUG_EXPR, the same
        # region-as-a-single-expression already used by page 15) instead of
        # one full-table query per region state.
        dealers_region_subgroups = self._safe('dealers_region_subgroups', lambda: self._breakdown_ty_ly_multi_topn(
            year, month, 'ytd', franchise, REGION_SLUG_EXPR, list(REGION_SUBQUERY.keys()),
            "bi_psgroupname", limit=8, extra_clauses=["bi_csttypecode = '01'", f"({mg_expr}) = 'Split'"],
            qty_gate=qty_gate),
            {slug: [] for slug in REGION_SUBQUERY})

        # Page 17: projects-productgroups (YTD).
        projects_productgroups = self._safe('projects_productgroups', lambda: self._breakdown_ty_ly(
            year, month, 'ytd', franchise, "bi_pgroupname", extra_clauses=["bi_csttypecode = '04'"], limit=8,
            qty_gate=qty_gate), [])

        bundle = {
            'period': {'year': year, 'month': month, 'monthName': MONTH_NAMES[month - 1],
                       'label': month_label, 'quarter': quarter, 'quarterLabel': f'Q{quarter} {year}'},
            'hasPrevYear': has_prev_year,
            'kpis': kpis,
            'totalCompany': total_company,
            'deptTrend': dept_trend,
            'regionTrend': region_trend,
            'quarterlyCompany': quarterly_company,
            'quarterlyDepartment': quarterly_department,
            'quarterlyRegion': quarterly_region,
            'mainGroups': main_groups,
            'subgroupsTop8': subgroups_top8,
            'departments': departments,
            'lcacSubgroups': lcac_subgroups,
            'top3ByChannel': top3_by_channel,
            'regionsByChannel': regions_by_channel,
            'dealersRegionMaingroups': dealers_region_maingroups,
            'dealersRegionSubgroups': dealers_region_subgroups,
            'projectsProductgroups': projects_productgroups,
        }

        narratives = {
            'total-company': self._narr_total_company(total_company['mtd'][0] if total_company['mtd'] else None,
                                                        total_company['ytd'][0] if total_company['ytd'] else None),
            'dept-region-trends': self._narr_breakdown(dept_trend, 'departments') + "\n" +
                                   self._narr_breakdown(region_trend, 'dealer regions'),
            'quarterly-company': self._narr_quarterly(quarterly_company, 'the company'),
        }
        for slug in CHANNEL_CODES:
            narratives[f'quarterly-department_{slug}'] = self._narr_quarterly(quarterly_department[slug], CHANNEL_LABELS[slug])
        for slug in REGION_SUBQUERY:
            narratives[f'quarterly-region_{slug}'] = self._narr_quarterly(quarterly_region[slug], REGION_LABELS[slug])
        narratives['main-groups'] = self._narr_breakdown(main_groups['ytd'], 'main groups')
        narratives['subgroups-top8'] = self._narr_breakdown(subgroups_top8, 'sub-groups')
        narratives['subgroups-distribution'] = self._narr_breakdown(subgroups_top8, 'sub-groups')
        narratives['maingroups-distribution'] = self._narr_breakdown(main_groups['ytd'], 'main groups')
        narratives['departments'] = self._narr_breakdown(departments['ytd'], 'departments')
        narratives['departments-distribution'] = self._narr_breakdown(departments['ytd'], 'departments')
        narratives['lcac-subgroups'] = self._narr_breakdown(lcac_subgroups['ytd'], 'LCAC sub-groups')
        for slug in CHANNEL_CODES:
            narratives[f'top3-by-channel_{slug}'] = self._narr_breakdown(top3_by_channel[slug], 'sub-groups', value_key='qty', budget_key='budgetQty', prev_key='prevYearQty')
            narratives[f'regions-by-channel_{slug}'] = self._narr_breakdown(regions_by_channel[slug], 'regions', value_key='qty', budget_key='budgetQty', prev_key='prevYearQty')
        for slug in REGION_SUBQUERY:
            narratives[f'dealers-region-maingroups_{slug}'] = self._narr_breakdown(dealers_region_maingroups[slug], 'main groups')
            narratives[f'dealers-region-subgroups_{slug}'] = self._narr_breakdown(dealers_region_subgroups[slug], 'sub-groups', value_key='qty', budget_key='budgetQty', prev_key='prevYearQty')
        narratives['projects-productgroups'] = self._narr_breakdown(projects_productgroups, 'product groups')

        # Combined-slide keys (no per-state suffix) for the 6 toggle pages,
        # which the PPTX export now renders as a single slide covering all
        # 4 states — see _combined_narrative and _build_pptx.
        narratives['quarterly-department'] = self._combined_narrative(narratives, 'quarterly-department', CHANNEL_LABELS)
        narratives['quarterly-region'] = self._combined_narrative(narratives, 'quarterly-region', REGION_LABELS)
        narratives['top3-by-channel'] = self._combined_narrative(narratives, 'top3-by-channel', CHANNEL_LABELS)
        narratives['regions-by-channel'] = self._combined_narrative(narratives, 'regions-by-channel', CHANNEL_LABELS)
        narratives['dealers-region-maingroups'] = self._combined_narrative(narratives, 'dealers-region-maingroups', REGION_LABELS)
        narratives['dealers-region-subgroups'] = self._combined_narrative(narratives, 'dealers-region-subgroups', REGION_LABELS)

        # Plain text is what _pptx_add_notes_sidebar expects (it marks
        # internally — see _build_pptx); the marked form is for on-screen
        # display and as the baseline import_notes diffs a re-imported
        # file's notes against. Marking twice (once here, again inside
        # _pptx_add_notes_sidebar) would corrupt the round-trip, so
        # bundle['narratives'] is ONLY ever marked once, right here.
        bundle['narrativesPlain'] = narratives
        bundle['narratives'] = {k: _pptx_marked_text(v) for k, v in narratives.items()}
        return bundle

    def _fetch_notes_map(self, period, franchise):
        notes = request.env['pbi.dashboard.note'].sudo().search([
            ('year', '=', period), ('franchise', '=', franchise),
            ('customer_type', '=', 'all'), ('region', '=', 'all'),
            ('key', 'like', 'sales_mail_%'),
        ])
        return {n.key[len('sales_mail_'):]: n.text for n in notes}

    # ------------------------------------------------------------------
    # route
    # ------------------------------------------------------------------
    def _sales_mail_data(self, has_access, period=None, franchise='Midea'):
        if not has_access:
            return {'error': 'You do not have access to this dashboard.'}
        try:
            year, month = (None, None)
            if period:
                try:
                    y, m = period.split('-')
                    year, month = int(y), int(m)
                except (ValueError, AttributeError):
                    year, month = None, None
            if year is None or month is None:
                year, month = self._latest_period()
            if year is None:
                return {'error': 'No sales data is available.'}

            bundle = self._fetch_bundle(year, month, franchise)
            bundle['notes'] = self._fetch_notes_map(f"{year}-{month:02d}", franchise)
            bundle['yearOptions'] = self._year_options()
            return bundle
        except Exception as e:
            return {'error': str(e)}

    @http.route('/pbi_dashboards/sales_mail/data', type='json', auth='user')
    def sales_mail_data(self, period=None, franchise='Midea'):
        return self._sales_mail_data(self._has_access(), period, franchise)

    @http.route('/pbi_dashboards/sales_mail_new/data', type='json', auth='user')
    def sales_mail_data_new(self, period=None, franchise='Midea'):
        return self._sales_mail_data(self._has_access_new(), period, franchise)

    # ------------------------------------------------------------------
    # PPTX export
    # ------------------------------------------------------------------
    def _pptx_style_series(self, chart, colors):
        for i, series in enumerate(chart.plots[0].series):
            if i < len(colors):
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = colors[i]

    def _pptx_style_pie_points(self, chart, colors):
        points = chart.plots[0].series[0].points
        for i, pt in enumerate(points):
            if i < len(colors):
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = colors[i]

    def _pptx_kpi_block(self, fmt_fn, this_year, target, last_year, prefix):
        """Mirrors renderKpis()'s `block()` helper in
        sales_mail_dashboard_new.js: This Year/Target/Last Year, each tile
        carrying the same vs-LY/Achv %-badge folded in (This Year -> vs LY,
        Target -> Achv, Last Year -> none)."""
        def sub(value, compared, badge_prefix):
            if not compared:
                return None, None
            p = value / compared * 100
            return f"{badge_prefix} {p:.1f}%", p >= 100
        ly_text, ly_good = sub(this_year, last_year, 'vs LY')
        tgt_text, tgt_good = sub(this_year, target, 'Achv')
        return [
            {'value': fmt_fn(this_year), 'label': f'{prefix} - This Year', 'sub_text': ly_text, 'sub_good': ly_good},
            {'value': fmt_fn(target), 'label': f'{prefix} - Target', 'sub_text': tgt_text, 'sub_good': tgt_good},
            {'value': fmt_fn(last_year), 'label': f'{prefix} - Last Year', 'sub_text': None, 'sub_good': None},
        ]

    def _pptx_add_kpi_tile(self, slide, left, top, width, height, tile, color):
        """One colored tile card — mirrors the web's `.pbi-kpi` (value/
        label/sub-badge stack on a solid-color rounded box)."""
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                        Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.shadow.inherit = False
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Pt(4)
        tf.margin_top = tf.margin_bottom = Pt(2)

        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run()
        r0.text = tile['value']
        r0.font.bold = True
        r0.font.size = Pt(15)
        r0.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

        p1 = tf.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = tile['label']
        r1.font.size = Pt(8.5)
        r1.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

        if tile['sub_text']:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = tile['sub_text']
            r2.font.bold = True
            r2.font.size = Pt(8)
            r2.font.color.rgb = PPTX_TILE_GOOD if tile['sub_good'] else PPTX_TILE_BAD

    def _pptx_add_tiles_slide(self, prs, layout, kpis, period_label, franchise, year):
        """Slide 1 = the dashboard's title/source line PLUS the on-screen
        12-tile KPI grid (Value + Qty, each MTD/YTD x This Year/Target/Last
        Year) rendered as actual colored tile cards — mirrors
        .pbi-kpi-row.six in sales_mail_dashboard_new.js, not a data table."""
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.55))
        p = tb.text_frame.paragraphs[0]
        p.text = "Sales Mail Content"
        p.font.size = Pt(28)
        p.font.bold = True
        franchise_label = "All Franchises" if franchise == "all" else franchise
        sub = slide.shapes.add_textbox(Inches(0.4), Inches(0.75), Inches(12.5), Inches(0.35))
        sub.text_frame.paragraphs[0].text = f"Source: bidata  ·  {franchise_label}  ·  {period_label}"
        sub.text_frame.paragraphs[0].font.size = Pt(13)

        def kpi_row(top, title, mtd, ytd, fmt_fn):
            htb = slide.shapes.add_textbox(Inches(PPTX_MARGIN), Inches(top), Inches(4), Inches(0.3))
            htb.text_frame.paragraphs[0].text = title
            htb.text_frame.paragraphs[0].font.bold = True
            htb.text_frame.paragraphs[0].font.size = Pt(15)
            tiles_top = top + 0.35
            tiles_h = 1.6
            tiles = self._pptx_kpi_block(fmt_fn, *mtd, 'MTD') + self._pptx_kpi_block(fmt_fn, *ytd, 'YTD')
            gap = 0.15
            w = (PPTX_FULL_W - (len(tiles) - 1) * gap) / len(tiles)
            for i, tile in enumerate(tiles):
                left = PPTX_MARGIN + i * (w + gap)
                self._pptx_add_kpi_tile(slide, left, tiles_top, w, tiles_h, tile, PPTX_TILE_COLORS[i % 3])
            return tiles_top + tiles_h

        next_top = kpi_row(1.3, "Value",
                            (kpis['mtdThisYear'], kpis['mtdTarget'], kpis['mtdLastYear']),
                            (kpis['ytdThisYear'], kpis['ytdTarget'], kpis['ytdLastYear']), _pfmt_m)
        kpi_row(next_top + 0.3, "Quantity",
                (kpis['mtdQtyThisYear'], kpis['mtdQtyTarget'], kpis['mtdQtyLastYear']),
                (kpis['ytdQtyThisYear'], kpis['ytdQtyTarget'], kpis['ytdQtyLastYear']), _pfmt)
        return slide

    def _pptx_grid_cells(self, n, stacked=False):
        """(left, top, width, height) in inches for n charts (1, 2, 4 or 8)
        within the chart area of a 13.333x7.5in slide — i.e. everything
        left of the notes sidebar (PPTX_NOTES_LEFT..). 1/2/4 match the web
        page's own .pbi-card-grid.two/.four layout (side-by-side / 2x2). 8
        is used only by the 2-chart-per-toggle-state pages consolidated
        onto one slide (4 states x 2 charts, arranged 4 columns x 2 rows).
        `stacked` mirrors the web's `.pbi-card-grid.two.stacked-full`
        modifier (grid-template-columns: 1fr) — full-width, one chart
        above the other — used by the n==2 pages whose XML carries that
        class (dept-region-trends, main-groups, subgroups-top8,
        lcac-subgroups, projects-productgroups) instead of the plain
        `.two` side-by-side pages (subgroups/maingroups/departments
        distribution donuts)."""
        left0, top0, total_w, total_h, gap = PPTX_MARGIN, PPTX_CONTENT_TOP, PPTX_CHART_AREA_W, PPTX_CONTENT_H, PPTX_GRID_GAP
        if n == 1:
            return [(left0, top0, total_w, total_h)]
        if n == 2:
            if stacked:
                h = (total_h - gap) / 2
                return [(left0, top0, total_w, h), (left0, top0 + h + gap, total_w, h)]
            w = (total_w - gap) / 2
            return [(left0, top0, w, total_h), (left0 + w + gap, top0, w, total_h)]
        if n == 4:
            w, h = (total_w - gap) / 2, (total_h - gap) / 2
            return [
                (left0, top0, w, h), (left0 + w + gap, top0, w, h),
                (left0, top0 + h + gap, w, h), (left0 + w + gap, top0 + h + gap, w, h),
            ]
        if n == 8:
            w, h = (total_w - 3 * gap) / 4, (total_h - gap) / 2
            return [
                (left0 + col * (w + gap), top0 + row * (h + gap), w, h)
                for row in range(2) for col in range(4)
            ]
        raise ValueError(f"unsupported chart count for a grid slide: {n}")

    def _pptx_add_bar_data_labels(self, chart, value_fmt):
        """Shows each bar/column's value above it, matching the web
        dashboard's `.bar-value` labels (groupedBarChart in
        sales_mail_dashboard_new.js) — native pptx charts render bare by
        default."""
        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.number_format = PPTX_AXIS_NUMBER_FORMAT[value_fmt]
        dl.number_format_is_linked = False
        dl.font.size = Pt(7)
        dl.position = XL_LABEL_POSITION.OUTSIDE_END

    def _pptx_add_pie_data_labels(self, chart):
        """Shows each slice's share %, matching the web dashboard's donut
        leader-line labels (donutChart in sales_mail_dashboard_new.js)."""
        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.show_percentage = True
        dl.show_value = False
        dl.number_format = '0.0%'
        dl.number_format_is_linked = False
        dl.font.size = Pt(8)

    def _pptx_add_combo_lines(self, chart, categories, lines):
        """Overlays achievementPct/yoyPct as a secondary-axis line series on
        top of an existing bar chart — mirrors the web's comboBarLineChart
        (Page 2: Department & Region Performance Trends), which the PPTX
        export previously dropped because python-pptx's high-level API has
        no combo/secondary-axis chart type. There's no way to build one
        through that API, so a `<c:lineChart>` plus its own hidden
        secondary c:catAx + visible right-hand c:valAx are spliced into the
        chart's plotArea as raw OOXML instead — python-pptx's chart.plots
        picks the new LinePlot straight back up afterward since it just
        re-scans the plotArea for known chart-type children, so no other
        code needs to know this series wasn't built the normal way."""
        plot_area = chart._chartSpace.chart.plotArea
        bar_chart_el = plot_area.find(qn('c:barChart'))
        base_idx = len(bar_chart_el.findall(qn('c:ser')))

        # Fixed sentinel axIds — only need to be unique within this one
        # chart's own plotArea, never across charts, so no need to derive
        # them from the bar chart's own (randomly-generated) axIds.
        sec_cat_ax_id, sec_val_ax_id = 500000001, 500000002

        def pt_xml(values):
            pts = "".join(f'<c:pt idx="{i}"><c:v>{v}</c:v></c:pt>' for i, v in enumerate(values) if v is not None)
            return f'<c:ptCount val="{len(values)}"/>{pts}'

        cat_xml = "".join(f'<c:pt idx="{i}"><c:v>{c}</c:v></c:pt>' for i, c in enumerate(categories))

        ser_xml = ""
        for i, line in enumerate(lines):
            idx = base_idx + i
            color_hex = "%02X%02X%02X" % (line['color'][0], line['color'][1], line['color'][2])
            ser_xml += f"""
            <c:ser>
              <c:idx val="{idx}"/>
              <c:order val="{idx}"/>
              <c:tx><c:strRef><c:f>Sheet1!$Z${i + 1}</c:f><c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>{line['label']}</c:v></c:pt></c:strCache></c:strRef></c:tx>
              <c:spPr><a:ln w="28575"><a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill></a:ln></c:spPr>
              <c:marker><c:symbol val="circle"/><c:size val="6"/><c:spPr><a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill><a:ln><a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill></a:ln></c:spPr></c:marker>
              <c:cat><c:strRef><c:f>Sheet1!$A$2:$A${len(categories) + 1}</c:f><c:strCache><c:ptCount val="{len(categories)}"/>{cat_xml}</c:strCache></c:strRef></c:cat>
              <c:val><c:numRef><c:f>Sheet1!$Y{i + 1}$2:$Y{i + 1}${len(categories) + 1}</c:f><c:numCache><c:formatCode>0.0&quot;%&quot;</c:formatCode>{pt_xml(line['values'])}</c:numCache></c:numRef></c:val>
              <c:smooth val="0"/>
            </c:ser>
            """

        line_chart_el = parse_xml(f"""<c:lineChart {nsdecls('c', 'a', 'r')}>
            <c:grouping val="standard"/>
            <c:varyColors val="0"/>
            {ser_xml}
            <c:marker val="1"/>
            <c:axId val="{sec_cat_ax_id}"/>
            <c:axId val="{sec_val_ax_id}"/>
          </c:lineChart>""")
        bar_chart_el.addnext(line_chart_el)

        sec_cat_ax_el = parse_xml(f"""<c:catAx {nsdecls('c', 'a', 'r')}>
            <c:axId val="{sec_cat_ax_id}"/>
            <c:scaling><c:orientation val="minMax"/></c:scaling>
            <c:delete val="1"/>
            <c:axPos val="b"/>
            <c:majorTickMark val="none"/>
            <c:minorTickMark val="none"/>
            <c:tickLblPos val="none"/>
            <c:crossAx val="{sec_val_ax_id}"/>
            <c:crosses val="autoZero"/>
            <c:auto val="1"/>
            <c:lblAlgn val="ctr"/>
            <c:lblOffset val="100"/>
            <c:noMultiLvlLbl val="0"/>
          </c:catAx>""")
        sec_val_ax_el = parse_xml(f"""<c:valAx {nsdecls('c', 'a', 'r')}>
            <c:axId val="{sec_val_ax_id}"/>
            <c:scaling><c:orientation val="minMax"/></c:scaling>
            <c:delete val="0"/>
            <c:axPos val="r"/>
            <c:numFmt formatCode="0&quot;%&quot;" sourceLinked="0"/>
            <c:majorTickMark val="out"/>
            <c:minorTickMark val="none"/>
            <c:tickLblPos val="nextTo"/>
            <c:crossAx val="{sec_cat_ax_id}"/>
            <c:crosses val="max"/>
          </c:valAx>""")
        # All chart-type elements (barChart, lineChart) must precede all
        # axis elements in c:plotArea's schema order — appending puts these
        # after the existing primary catAx/valAx, which is exactly right.
        plot_area.append(sec_cat_ax_el)
        plot_area.append(sec_val_ax_el)

    def _pptx_add_notes_sidebar(self, slide, notes):
        """Visible notes column to the right of the chart grid, mirroring
        the web page's `.chart-row` layout (chart-col + .narrative sidebar,
        grid-template-columns: 1fr 340px). This is the ONLY place notes
        live — PowerPoint's separate Notes pane is left empty on purpose
        (see _pptx_add_multi_chart_slide); _pptx_find_notes_sidebar locates
        this exact textbox again on re-import by its fixed left position."""
        box = slide.shapes.add_textbox(Inches(PPTX_NOTES_LEFT), Inches(PPTX_CONTENT_TOP),
                                        Inches(PPTX_NOTES_W), Inches(PPTX_CONTENT_H))
        box.text_frame.word_wrap = True
        # A combined toggle-page narrative (4 states concatenated, see
        # _combined_narrative) can run long enough to overflow the fixed
        # sidebar height — shrink the font in PowerPoint rather than let it
        # spill past the slide edge.
        box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        if notes:
            _pptx_fill_marked_text(box.text_frame, notes, heading_size=Pt(12), body_size=Pt(9))

    def _pptx_add_multi_chart_slide(self, prs, layout, title, notes, specs, stacked=False):
        """One slide per web PAGE (or per toggle state), with that page's
        1/2/4 charts arranged in a grid to the left and its narrative in a
        visible sidebar to the right — mirrors the on-screen `.chart-row`
        layout. `specs` is a list of dicts from _bar_spec/_pie_spec, each
        either {'subtitle','type','categories','series','colors','is_pie'}
        or {'subtitle','empty': True} when that chart's data is empty.
        `stacked` is forwarded to _pptx_grid_cells for the n==2 pages that
        need chart-above-chart instead of side-by-side (see its docstring)."""
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.55))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True

        for spec, (left, top, width, height) in zip(specs, self._pptx_grid_cells(len(specs), stacked=stacked)):
            subtitle = spec.get('subtitle')
            chart_top, chart_height = top, height
            if subtitle:
                stb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.28))
                sp = stb.text_frame.paragraphs[0]
                sp.text = subtitle
                sp.font.size = Pt(12)
                sp.font.bold = True
                chart_top, chart_height = top + 0.3, height - 0.3
            if spec.get('empty'):
                ntb = slide.shapes.add_textbox(Inches(left), Inches(chart_top), Inches(width), Inches(0.4))
                ntb.text_frame.paragraphs[0].text = "No data"
                ntb.text_frame.paragraphs[0].font.size = Pt(11)
                continue
            chart_data = CategoryChartData()
            chart_data.categories = spec['categories']
            for name, values in spec['series']:
                chart_data.add_series(name, values)
            gframe = slide.shapes.add_chart(spec['type'], Inches(left), Inches(chart_top),
                                             Inches(width), Inches(chart_height), chart_data)
            chart = gframe.chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.font.size = Pt(8)
            if spec.get('is_pie'):
                self._pptx_style_pie_points(chart, spec['colors'])
                self._pptx_add_pie_data_labels(chart)
            else:
                self._pptx_style_series(chart, spec['colors'])
                # Native pptx defaults to overlap=0 (bars in the same
                # cluster touch edge-to-edge) — the web's groupedBarChart
                # always leaves a small gap between them (its `gap = 6`
                # out of up to `barW = 32`), so pull them apart to match.
                chart.plots[0].overlap = -30
                # Native pptx charts default the value axis to "General"
                # format (raw numbers like 630643140) — match the web
                # dashboard's M/K-abbreviated axis labels instead.
                chart.value_axis.tick_labels.number_format = PPTX_AXIS_NUMBER_FORMAT[spec.get('value_fmt', 'M')]
                chart.value_axis.tick_labels.number_format_is_linked = False
                self._pptx_add_bar_data_labels(chart, spec.get('value_fmt', 'M'))
                if spec.get('lines'):
                    self._pptx_add_combo_lines(chart, spec['categories'], spec['lines'])

        # Notes live ONLY in the visible sidebar textbox now — PowerPoint's
        # separate Notes pane (shown at the bottom of the editing window)
        # is intentionally left untouched/empty; _import_notes reads the
        # sidebar textbox back via _pptx_find_notes_sidebar instead of
        # slide.notes_slide.
        self._pptx_add_notes_sidebar(slide, notes)
        return slide

    def _qty_sorted(self, rows):
        """Every breakdown list from _fetch_bundle is ordered by This-Year
        SALES descending (shared with the paired Value chart/slide). A Qty
        chart/slide needs its OWN This-Year-Qty-descending order instead —
        sort a copy here so the Value chart's list/order is untouched."""
        return sorted(rows or [], key=lambda r: r.get('qty') or 0, reverse=True)

    def _bar_spec(self, rows, series_defs, subtitle=None):
        if not rows:
            return {'subtitle': subtitle, 'empty': True}
        cats = [r['label'] or '—' for r in rows]
        series = [(label, [r.get(key) or 0 for r in rows]) for key, label in series_defs]
        # series_defs is always value_series (sales/budget/prevYearSales,
        # web's fmtM) or qty_series (qty/budgetQty/prevYearQty, web's
        # fmtK) — detect which from the first series key so every call
        # site gets the matching axis format without threading a new arg.
        value_fmt = 'K' if series_defs[0][0] in ('qty', 'budgetQty', 'prevYearQty') else 'M'
        return {'subtitle': subtitle, 'type': XL_CHART_TYPE.COLUMN_CLUSTERED, 'categories': cats, 'series': series,
                'colors': [PPTX_COLOR_2025, PPTX_COLOR_BUDGET, PPTX_COLOR_2024], 'is_pie': False, 'value_fmt': value_fmt}

    def _combo_spec(self, rows, series_defs, line_defs, subtitle=None):
        """_bar_spec plus a `lines` key carrying achievementPct/yoyPct as a
        secondary-axis overlay — matches the web's comboBarLineChart, used
        only by Page 2 (Department & Region Performance Trends). Consumed
        by _pptx_add_multi_chart_slide -> _pptx_add_combo_lines.
        `line_defs` is a list of (row_key, label, color) tuples."""
        spec = self._bar_spec(rows, series_defs, subtitle)
        if spec.get('empty'):
            return spec
        spec['lines'] = [
            {'label': label, 'color': color, 'values': [r.get(key) for r in rows]}
            for key, label, color in line_defs
        ]
        return spec

    def _pie_spec(self, rows, value_key, subtitle=None):
        rows = [r for r in (rows or []) if (r.get(value_key) or 0) > 0 and r.get('label') != 'Total']
        if not rows:
            return {'subtitle': subtitle, 'empty': True}
        cats = [r['label'] for r in rows]
        series = [('Share', [r.get(value_key) or 0 for r in rows])]
        return {'subtitle': subtitle, 'type': XL_CHART_TYPE.PIE, 'categories': cats, 'series': series,
                'colors': PPTX_PIE_PALETTE, 'is_pie': True}

    def _build_pptx(self, year, month, franchise):
        bundle = self._fetch_bundle(year, month, franchise)
        # Plain text — _pptx_add_multi_chart_slide -> _pptx_add_notes_sidebar
        # marks it internally. Using the already-marked bundle['narratives']
        # here would mark twice and corrupt the import round-trip.
        narratives = bundle['narrativesPlain']
        period_label = bundle['period']['label']

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        self._pptx_add_tiles_slide(prs, blank, bundle['kpis'], period_label, franchise, year)

        def page_slide(title, notes_key, specs, stacked=False):
            self._pptx_add_multi_chart_slide(prs, blank, title, narratives.get(notes_key, ''), specs, stacked=stacked)

        bar = self._bar_spec
        combo = self._combo_spec
        pie = self._pie_spec
        qsort = self._qty_sorted
        value_series = [('sales', 'This Year'), ('budget', 'Target'), ('prevYearSales', 'Last Year')]
        qty_series = [('qty', 'This Year'), ('budgetQty', 'Target'), ('prevYearQty', 'Last Year')]
        trend_lines = [('achievementPct', 'Vs. Target', PPTX_LINE_ACHIEVEMENT), ('yoyPct', 'Vs. Last Year', PPTX_LINE_YOY)]

        # Page 1: Total Company — 4 charts, 2x2 (matches .pbi-card-grid.two x2 rows on screen)
        tc = bundle['totalCompany']
        page_slide(f"1. Total Company ({period_label})", 'total-company', [
            bar(tc['mtd'], value_series, f"{period_label} - Value"),
            bar(qsort(tc['mtd']), qty_series, f"{period_label} - Quantity"),
            bar(tc['ytd'], value_series, f"{year} YTD - Value"),
            bar(qsort(tc['ytd']), qty_series, f"{year} YTD - Quantity"),
        ])

        # Page 2: Department & Region Performance Trends — 2 charts stacked
        # (one above the other, matches .pbi-card-grid.two.stacked-full on
        # screen), each with the achievement%/YoY% lines overlaid on a
        # secondary axis via _combo_spec/_pptx_add_combo_lines — mirrors
        # the web's comboBarLineChart.
        page_slide(f"2. Department & Region Performance Trends ({period_label})", 'dept-region-trends', [
            combo(bundle['deptTrend'], value_series, trend_lines, "By Department"),
            combo(bundle['regionTrend'], value_series, trend_lines, "By Dealer Region"),
        ], stacked=True)

        # Page 3: Quarterly Sales Progression — Company — 1 chart, full slide
        page_slide(f"3. Quarterly Sales Progression — Company ({year})", 'quarterly-company', [
            bar(bundle['quarterlyCompany'], value_series),
        ])

        # Page 4: Quarterly Progression by Department — one slide PER TAB
        # (matches the PDF's per-tab-print pages), each with its own
        # per-state narrative instead of the old single 2x2-grid slide with
        # a combined narrative.
        for slug, label in CHANNEL_LABELS.items():
            page_slide(f"4. Quarterly Progression by Department — {label} ({year})",
                       f'quarterly-department_{slug}',
                       [bar(bundle['quarterlyDepartment'][slug], value_series, label)])

        # Page 5: Quarterly Progression by Region — one slide per tab.
        for slug, label in REGION_LABELS.items():
            page_slide(f"5. Quarterly Progression by Region — {label} ({year})",
                       f'quarterly-region_{slug}',
                       [bar(bundle['quarterlyRegion'][slug], value_series, label)])

        # Page 6: Main Groups — MTD & YTD stacked (.stacked-full on screen)
        page_slide(f"6. Main Groups ({period_label})", 'main-groups', [
            bar(bundle['mainGroups']['mtd'], value_series, period_label),
            bar(bundle['mainGroups']['ytd'], value_series, f"{year} YTD"),
        ], stacked=True)

        # Page 7: Product Sub-Groups — Top 8 (YTD) — Value/Qty stacked (.stacked-full on screen)
        page_slide(f"7. Product Sub-Groups — Top 8 (YTD {period_label})", 'subgroups-top8', [
            bar(bundle['subgroupsTop8'], value_series, "Value"),
            bar(qsort(bundle['subgroupsTop8']), qty_series, "Quantity"),
        ], stacked=True)

        # Page 8: Sub-Groups Distribution — LY vs TY side by side
        page_slide(f"8. Sub-Groups Distribution — LY vs TY (Top 8, YTD {period_label})", 'subgroups-distribution', [
            pie(bundle['subgroupsTop8'], 'prevYearSales', "Last Year"),
            pie(bundle['subgroupsTop8'], 'sales', "This Year"),
        ])

        # Page 9: Main Groups Distribution — LY vs TY side by side
        page_slide(f"9. Main Groups Distribution — LY vs TY (YTD {period_label})", 'maingroups-distribution', [
            pie(bundle['mainGroups']['ytd'], 'prevYearSales', "Last Year"),
            pie(bundle['mainGroups']['ytd'], 'sales', "This Year"),
        ])

        # Page 10: Departments — MTD & YTD — 4 charts, 2x2
        dpt = bundle['departments']
        page_slide(f"10. Departments ({period_label})", 'departments', [
            bar(dpt['mtd'], value_series, f"{period_label} - Value"),
            bar(qsort(dpt['mtd']), qty_series, f"{period_label} - Quantity"),
            bar(dpt['ytd'], value_series, f"{year} YTD - Value"),
            bar(qsort(dpt['ytd']), qty_series, f"{year} YTD - Quantity"),
        ])

        # Page 11: Departments Distribution — LY vs TY — 4 donuts, 2x2
        page_slide(f"11. Departments Distribution — LY vs TY (YTD {period_label})", 'departments-distribution', [
            pie(dpt['ytd'], 'prevYearSales', "Value — Last Year"),
            pie(dpt['ytd'], 'sales', "Value — This Year"),
            pie(qsort(dpt['ytd']), 'prevYearQty', "Qty — Last Year"),
            pie(qsort(dpt['ytd']), 'qty', "Qty — This Year"),
        ])

        # Page 12: LCAC Sub-Groups — MTD & YTD stacked (.stacked-full on screen)
        page_slide(f"12. LCAC Sub-Groups ({period_label})", 'lcac-subgroups', [
            bar(bundle['lcacSubgroups']['mtd'], value_series, period_label),
            bar(bundle['lcacSubgroups']['ytd'], value_series, f"{year} YTD"),
        ], stacked=True)

        # Page 13: Top-3 Sub-Groups by Channel — one slide PER TAB (matches
        # the PDF's per-tab-print pages), Value stacked above Quantity
        # (matches the on-screen .stacked-full for this pair) instead of the
        # old single 4x2-grid slide with a combined narrative.
        for slug, label in CHANNEL_LABELS.items():
            rows = bundle['top3ByChannel'][slug]
            page_slide(f"13. Top-3 Sub-Groups by Channel (YTD {period_label}) — {label}",
                       f'top3-by-channel_{slug}',
                       [bar(rows, value_series, "Value"), bar(qsort(rows), qty_series, "Quantity")],
                       stacked=True)

        # Page 14: Regions by Channel — same one-slide-per-tab treatment.
        for slug, label in CHANNEL_LABELS.items():
            rows = bundle['regionsByChannel'][slug]
            page_slide(f"14. Regions by Channel (YTD {period_label}) — {label}",
                       f'regions-by-channel_{slug}',
                       [bar(rows, value_series, "Value"), bar(qsort(rows), qty_series, "Quantity")],
                       stacked=True)

        # Page 15: Dealers by Region — Main Groups — one slide per tab.
        for slug, label in REGION_LABELS.items():
            page_slide(f"15. Dealers by Region — Main Groups (YTD {period_label}) — {label}",
                       f'dealers-region-maingroups_{slug}',
                       [bar(bundle['dealersRegionMaingroups'][slug], value_series, label)])

        # Page 16: Dealers by Region — Split Sub-Groups — one slide per tab.
        for slug, label in REGION_LABELS.items():
            page_slide(f"16. Dealers by Region — Split Sub-Groups (YTD {period_label}) — {label}",
                       f'dealers-region-subgroups_{slug}',
                       [bar(qsort(bundle['dealersRegionSubgroups'][slug]), qty_series, label)])

        # Page 17: Projects — Product Groups — Value/Qty stacked (.stacked-full on screen)
        page_slide(f"17. Projects — Product Groups (YTD {period_label})", 'projects-productgroups', [
            bar(bundle['projectsProductgroups'], value_series, "Value"),
            bar(qsort(bundle['projectsProductgroups']), qty_series, "Quantity"),
        ], stacked=True)

        prs.core_properties.keywords = f"pbi_dashboards_mail|period={year}-{month:02d}|franchise={franchise}"
        stream = BytesIO()
        prs.save(stream)
        return stream.getvalue()

    def _export_pptx(self, has_access, period=None, franchise='Midea'):
        if not has_access:
            return request.not_found()
        year, month = (None, None)
        if period:
            try:
                y, m = period.split('-')
                year, month = int(y), int(m)
            except (ValueError, AttributeError):
                year, month = None, None
        if year is None or month is None:
            year, month = self._latest_period()
        if year is None:
            return request.make_response("No sales data is available.", status=404)
        try:
            data = self._build_pptx(year, month, franchise)
        except Exception as e:
            return request.make_response(f"Error building PPTX: {e}", status=500)
        headers = [
            ('Content-Type', 'application/vnd.openxmlformats-officedocument.presentationml.presentation'),
            ('Content-Disposition', 'attachment; filename="Sales_Mail_Content.pptx"'),
            ('Content-Length', len(data)),
        ]
        return request.make_response(data, headers=headers)

    @http.route('/pbi_dashboards/sales_mail/export.pptx', type='http', auth='user')
    def export_pptx(self, period=None, franchise='Midea', **kwargs):
        return self._export_pptx(self._has_access(), period, franchise)

    @http.route('/pbi_dashboards/sales_mail_new/export.pptx', type='http', auth='user')
    def export_pptx_new(self, period=None, franchise='Midea', **kwargs):
        return self._export_pptx(self._has_access_new(), period, franchise)

    # ------------------------------------------------------------------
    # PPTX notes import
    # ------------------------------------------------------------------
    def _import_notes(self, has_access, period=None, franchise='Midea'):
        def _json(payload, status=200):
            return request.make_response(json.dumps(payload), headers=[('Content-Type', 'application/json')], status=status)

        if not has_access:
            return _json({'error': 'You do not have access to this dashboard.'}, status=403)

        upload = request.httprequest.files.get('file')
        if not upload:
            return _json({'error': 'No file uploaded.'}, status=400)
        try:
            prs = Presentation(BytesIO(upload.read()))
        except Exception as e:
            return _json({'error': f'Could not read PowerPoint file: {e}'}, status=400)

        keywords = prs.core_properties.keywords or ''
        m = re.match(r'^pbi_dashboards_mail\|period=(?P<period>[^|]*)\|franchise=(?P<franchise>.*)$', keywords)
        if m:
            period = m.group('period')
            franchise = m.group('franchise')
        if not period:
            return _json({'error': 'Could not determine which period this file belongs to.'}, status=400)
        year, mth = int(period.split('-')[0]), int(period.split('-')[1])

        notes_map = {}
        for slide in prs.slides:
            title = _pptx_slide_title(slide)
            key = _pptx_note_section_key(title)
            if not key:
                continue
            sidebar = _pptx_find_notes_sidebar(slide)
            if not sidebar:
                continue
            text = _pptx_read_notes_marked(sidebar.text_frame).strip()
            if text:
                notes_map[key] = text

        try:
            baseline = dict(self._fetch_bundle(year, mth, franchise)['narratives'])
        except Exception:
            baseline = {}

        Note = request.env['pbi.dashboard.note'].sudo()
        applied = {}
        for key, text in notes_map.items():
            db_key = f'sales_mail_{key}'
            existing = Note.search([
                ('key', '=', db_key), ('year', '=', period), ('franchise', '=', franchise),
                ('customer_type', '=', 'all'), ('region', '=', 'all'),
            ], limit=1)
            if text.strip() == (baseline.get(key) or '').strip():
                if existing:
                    existing.unlink()
                continue
            applied[key] = text
            if existing:
                existing.text = text
            else:
                Note.create({'key': db_key, 'year': period, 'franchise': franchise,
                              'customer_type': 'all', 'region': 'all', 'text': text})

        return _json({'period': period, 'franchise': franchise, 'notes': applied})

    @http.route('/pbi_dashboards/sales_mail/import_notes', type='http', auth='user', methods=['POST'])
    def import_notes(self, period=None, franchise='Midea', **kwargs):
        return self._import_notes(self._has_access(), period, franchise)

    @http.route('/pbi_dashboards/sales_mail_new/import_notes', type='http', auth='user', methods=['POST'])
    def import_notes_new(self, period=None, franchise='Midea', **kwargs):
        return self._import_notes(self._has_access_new(), period, franchise)
